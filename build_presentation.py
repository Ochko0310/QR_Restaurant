"""Build a ~8-10 min defense presentation for the QR Restaurant system."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

PRIMARY = RGBColor(0xF9, 0x73, 0x16)      # orange
DARK = RGBColor(0x0F, 0x17, 0x23)         # near-black
CARD = RGBColor(0x1B, 0x24, 0x33)         # dark slate
MUTED = RGBColor(0x94, 0xA3, 0xB8)        # slate-400
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x34, 0xD3, 0x99)       # emerald

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_bg(slide, color=DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def add_text(slide, left, top, width, height, text, *, size=18,
             bold=False, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.55), color=PRIMARY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_card(slide, left, top, width, height, *, fill=CARD, border=PRIMARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.25)
    return box


def header(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.6), Inches(0.55))
    add_text(slide, Inches(0.85), Inches(0.45), Inches(12), Inches(0.75),
             title, size=30, bold=True)
    if subtitle:
        add_text(slide, Inches(0.85), Inches(1.05), Inches(12), Inches(0.45),
                 subtitle, size=14, color=MUTED)
    # footer line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55),
                                  Inches(12.15), Emu(9525))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY


def footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.35),
             "QR Ширээ захиалгын систем  •  Төгсөлтийн төсөл", size=10, color=MUTED)
    add_text(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
             f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Title ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
# left decorative stripe
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SH)
stripe.line.fill.background()
stripe.fill.solid()
stripe.fill.fore_color.rgb = PRIMARY

add_text(s, Inches(1.0), Inches(1.6), Inches(11.5), Inches(0.5),
         "ТӨГСӨЛТИЙН АЖЛЫН ТАНИЛЦУУЛГА", size=16, bold=True, color=PRIMARY)
add_text(s, Inches(1.0), Inches(2.15), Inches(11.5), Inches(1.4),
         "QR код дээр суурилсан\nресторан ширээ захиалгын систем",
         size=44, bold=True)
add_text(s, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.5),
         "Бодит цагийн захиалга · Нөөцийн удирдлага · Олон зочны дэмжлэг",
         size=18, color=MUTED)

# author box
box = add_card(s, Inches(1.0), Inches(5.5), Inches(6.5), Inches(1.35))
add_text(s, Inches(1.25), Inches(5.62), Inches(6), Inches(0.4),
         "Боловсруулсан", size=11, color=MUTED)
add_text(s, Inches(1.25), Inches(5.9), Inches(6), Inches(0.5),
         "Оюутны нэр", size=20, bold=True)
add_text(s, Inches(1.25), Inches(6.3), Inches(6), Inches(0.4),
         "Удирдагч багш: ____________", size=12, color=MUTED)

add_text(s, Inches(8.5), Inches(5.62), Inches(4), Inches(0.4),
         "Огноо", size=11, color=MUTED)
add_text(s, Inches(8.5), Inches(5.9), Inches(4), Inches(0.5),
         "2026 он", size=20, bold=True)

# ─── Slide 2: Problem ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Асуудлын тодорхойлолт", "Одоогийн рестораны үйлчилгээний сул талууд")

problems = [
    ("Хүлээлт их", "Зочин цэс, төлбөр, үйлчлэгчийг удаан хүлээж, сэтгэл ханамж буурдаг."),
    ("Хүний алдаа", "Гар захиалга буруу бичигдэх, давхардах, хасалт хийгдэхгүй байх."),
    ("Нөөц удирдлагагүй", "Ундаа, дарс зэрэг гаднаас орж ирсэн бараа хяналтгүй дуусдаг."),
    ("Бодит цаг алга", "Менежер борлуулалт, ачаалалтай цагийг шууд харж чадахгүй."),
]
cols = 2
gap = Inches(0.3)
card_w = (Inches(12.15) - gap) / cols
card_h = Inches(2.3)
start_top = Inches(2.05)
for i, (title, body) in enumerate(problems):
    r, c = divmod(i, cols)
    left = Inches(0.6) + c * (card_w + gap)
    top = start_top + r * (card_h + Inches(0.25))
    add_card(s, left, top, card_w, card_h)
    add_text(s, left + Inches(0.25), top + Inches(0.2), card_w - Inches(0.4),
             Inches(0.5), title, size=20, bold=True, color=PRIMARY)
    add_text(s, left + Inches(0.25), top + Inches(0.85), card_w - Inches(0.4),
             Inches(1.3), body, size=14, color=WHITE)
footer(s, 2, 12)

# ─── Slide 3: Objectives ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Зорилго ба зорилтууд")

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Ерөнхий зорилго", size=18, bold=True, color=PRIMARY)
add_text(s, Inches(0.6), Inches(2.35), Inches(12), Inches(0.9),
         "Зочин өөрийн утсаар QR код уншуулан цэс харж, захиалга өгөх;\n"
         "ажилтнууд бодит цагт захиалга, ширээ, нөөцийг удирдах нэгдмэл систем бий болгох.",
         size=15, color=WHITE)

add_text(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.5),
         "Нарийвчилсан зорилтууд", size=18, bold=True, color=PRIMARY)

objectives = [
    "Ширээ бүрт QR token үүсгэж, олон зочны зэрэгцэн захиалга өгөх боломж",
    "Захиалгын төлөв (pending → preparing → served → paid) автомат удирдлага",
    "Менежер, кассчин, үйлчлэгч, тогоочийн эрхийн түвшинтэй самбар",
    "Гаднаас орж ирсэн барааны нөөц, автомат хасалт, доод хэмжээний мэдэгдэл",
    "Бодит цагийн socket холболтоор бүх төхөөрөмж шууд шинэчлэгдэх",
    "Борлуулалт, топ хоол, ачаалалтай цагийн тайлан",
]
for i, obj in enumerate(objectives):
    top = Inches(4.35) + i * Inches(0.42)
    # bullet dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), top + Inches(0.12),
                             Inches(0.14), Inches(0.14))
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = PRIMARY
    add_text(s, Inches(1.05), top, Inches(11.8), Inches(0.4), obj, size=14)
footer(s, 3, 12)

# ─── Slide 4: Architecture ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Системийн архитектур", "Монорепо дээр суурилсан клиент-сервер бүтэц")

# Draw three tiers
tiers = [
    ("Клиент талууд", [
        "Зочны веб (QR → цэс → захиалга)",
        "Ажилтны самбар (менежер, кассчин,\nүйлчлэгч, тогооч)",
    ], RGBColor(0x2A, 0x37, 0x4B)),
    ("API сервер", [
        "Express 5 + TypeScript",
        "JWT + роль шалгалт",
        "Socket.IO (бодит цаг)",
        "REST: /api/orders, /api/menu,\n/api/inventory, /api/notifications",
    ], RGBColor(0x3B, 0x2D, 0x1C)),
    ("Өгөгдлийн сан", [
        "PostgreSQL",
        "Drizzle ORM + drizzle-kit push",
        "tables, orders, menu_items,\ninventory_items, notifications",
    ], RGBColor(0x1C, 0x33, 0x2C)),
]
col_w = Inches(4.0)
gap = Inches(0.25)
start_left = Inches(0.6)
for i, (title, lines, fill) in enumerate(tiers):
    left = start_left + i * (col_w + gap)
    top = Inches(2.0)
    height = Inches(4.8)
    add_card(s, left, top, col_w, height, fill=fill)
    add_text(s, left + Inches(0.3), top + Inches(0.25), col_w - Inches(0.6),
             Inches(0.55), title, size=20, bold=True, color=PRIMARY)
    for j, line in enumerate(lines):
        add_text(s, left + Inches(0.3), top + Inches(0.95) + j * Inches(0.9),
                 col_w - Inches(0.6), Inches(0.9), "• " + line, size=13)
# arrows between tiers
for i in range(2):
    ax = start_left + (i + 1) * col_w + i * gap - Emu(5000)
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, Inches(4.2),
                             gap + Emu(10000), Inches(0.4))
    arr.line.fill.background()
    arr.fill.solid()
    arr.fill.fore_color.rgb = PRIMARY
footer(s, 4, 12)

# ─── Slide 5: Tech stack ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Ашигласан технологи")

groups = [
    ("Frontend", ["React 18 + Vite", "TanStack Query", "Tailwind + shadcn/ui",
                  "Zustand (store)", "Recharts (график)", "Socket.IO client"]),
    ("Backend", ["Node.js + Express 5", "TypeScript", "Drizzle ORM",
                 "Socket.IO server", "JWT + bcrypt", "Multer (зураг)"]),
    ("Infra & хэрэгсэл", ["PostgreSQL", "pnpm monorepo", "esbuild",
                          "Orval (OpenAPI → React hook)", "LaTeX (thesis)",
                          "Git + GitHub"]),
]
col_w = Inches(4.0)
gap = Inches(0.25)
for i, (title, items) in enumerate(groups):
    left = Inches(0.6) + i * (col_w + gap)
    top = Inches(2.05)
    add_card(s, left, top, col_w, Inches(4.7))
    add_text(s, left + Inches(0.3), top + Inches(0.25), col_w - Inches(0.6),
             Inches(0.55), title, size=20, bold=True, color=PRIMARY)
    for j, item in enumerate(items):
        add_text(s, left + Inches(0.3), top + Inches(0.95) + j * Inches(0.55),
                 col_w - Inches(0.6), Inches(0.5), "▸ " + item, size=14)
footer(s, 5, 12)

# ─── Slide 6: QR ordering flow ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "QR захиалгын урсгал", "Зочин → захиалга → тогооч → төлбөр")

steps = [
    ("1", "QR skan", "Зочин ширээний QR-г уншуулж,\nцэс рүү шууд нэвтэрнэ"),
    ("2", "Session нээх", "Үйлчлэгч ширээг идэвхжүүлж\ntable_sessions үүсгэнэ"),
    ("3", "Захиалга", "Зочин хоол сонгож\nPOST /api/orders илгээнэ"),
    ("4", "Тогооч", "pending → preparing\nстатус шилжүүлнэ"),
    ("5", "Үйлчилгээ", "ready → served\n(бодит цагт сошиол)"),
    ("6", "Төлбөр", "Кассчин paid болгоход\nширээ автомат чөлөөлөгдөнө"),
]
step_w = Inches(2.0)
gap_x = Inches(0.1)
total_w = 6 * step_w + 5 * gap_x
start_x = (SW - total_w) / 2
row_top = Inches(3.0)
for i, (num, title, desc) in enumerate(steps):
    left = start_x + i * (step_w + gap_x)
    card = add_card(s, left, row_top, step_w, Inches(2.6))
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.7),
                              row_top + Inches(0.2), Inches(0.6), Inches(0.6))
    circ.line.fill.background()
    circ.fill.solid()
    circ.fill.fore_color.rgb = PRIMARY
    add_text(s, left + Inches(0.7), row_top + Inches(0.23), Inches(0.6),
             Inches(0.55), num, size=20, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.1), row_top + Inches(0.95), step_w - Inches(0.2),
             Inches(0.5), title, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.12), row_top + Inches(1.45),
             step_w - Inches(0.24), Inches(1.1), desc, size=11,
             color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.5),
         "Онцлог: нэг идэвхтэй session дээр олон гар утас зэрэг захиалга өгөх боломжтой",
         size=13, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s, 6, 12)

# ─── Slide 7: Features ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Үндсэн боломжууд")

features = [
    ("Олон зочин · нэг QR",
     "Ширээн дэх хүн болгон өөрийн утсаар\nзахиалга өгнө. customerId-аар тусгаарлана."),
    ("Бодит цагт шинэчлэл",
     "Socket.IO room-аар шинэ захиалга,\nстатус өөрчлөлт шууд харагдана."),
    ("Нөөцийн удирдлага",
     "Бараа нэмэхэд цэсэнд автомат нэмэгдэнэ.\nЗахиалга хасалт автомат."),
    ("Threshold мэдэгдэл",
     "Доод хэмжээг давж доошлоход\nхонхоор + DB-д хадгалагдана."),
    ("Төлбөр + ширээний нэгтгэл",
     "Нэг ширээний бүх идэвхтэй захиалгын\nнийт дүнг картан дээр харуулна."),
    ("Тайлан",
     "Борлуулалт, топ хоол, ачаалалтай цаг,\nтөлбөрийн арга — огноогоор шүүнэ."),
]
cols = 3
gap = Inches(0.22)
card_w = (Inches(12.15) - gap * (cols - 1)) / cols
card_h = Inches(2.3)
for i, (t, d) in enumerate(features):
    r, c = divmod(i, cols)
    left = Inches(0.6) + c * (card_w + gap)
    top = Inches(2.05) + r * (card_h + Inches(0.25))
    add_card(s, left, top, card_w, card_h)
    add_text(s, left + Inches(0.25), top + Inches(0.25), card_w - Inches(0.5),
             Inches(0.55), t, size=16, bold=True, color=PRIMARY)
    add_text(s, left + Inches(0.25), top + Inches(0.85), card_w - Inches(0.5),
             Inches(1.3), d, size=12, color=WHITE)
footer(s, 7, 12)

# ─── Slide 8: Staff dashboard roles ───────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Ажилтны самбар", "Роль тус бүрийн харах боломж ялгаатай")

roles = [
    ("Менежер", PRIMARY, [
        "Бүх цэс, ширээ, тайлан, зар, сэтгэгдэл",
        "Хоол, ангилал удирдлага",
        "Бараа, мэдэгдэл харах",
        "Борлуулалтын тайлан, топ хоол",
    ]),
    ("Кассчин", ACCENT, [
        "Захиалга, нэгтгэл, нөөц",
        "Төлбөр авах, paid төлөв",
        "Хоол void хийх эрх",
        "Бараа нэмэх, засах",
    ]),
    ("Үйлчлэгч", RGBColor(0x60, 0xA5, 0xFA), [
        "Идэвхтэй захиалгууд",
        "Ширээ идэвхжүүлэх (session)",
        "served төлөвт шилжүүлэх",
        "Урьдчилсан захиалга",
    ]),
    ("Тогооч", RGBColor(0xF4, 0x72, 0xB6), [
        "Preparing дараалал",
        "Ready болгон тэмдэглэх",
        "Зөвхөн идэвхтэй захиалга",
        "Хоол бүрийн тэмдэглэл",
    ]),
]
cols = 4
gap = Inches(0.15)
card_w = (Inches(12.15) - gap * (cols - 1)) / cols
card_h = Inches(4.6)
for i, (name, color, items) in enumerate(roles):
    left = Inches(0.6) + i * (card_w + gap)
    top = Inches(2.05)
    add_card(s, left, top, card_w, card_h, border=color)
    # colored header band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                              card_w, Inches(0.75))
    band.adjustments[0] = 0.22
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color
    add_text(s, left, top + Inches(0.13), card_w, Inches(0.5), name,
             size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    for j, it in enumerate(items):
        add_text(s, left + Inches(0.2), top + Inches(1.1) + j * Inches(0.75),
                 card_w - Inches(0.4), Inches(0.7), "• " + it, size=12)
footer(s, 8, 12)

# ─── Slide 9: Inventory ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Нөөцийн удирдлагын модуль")

add_text(s, Inches(0.6), Inches(1.9), Inches(7.5), Inches(0.5),
         "Ажлын урсгал", size=18, bold=True, color=PRIMARY)
flow = [
    "Менежер/кассчин бараа нэмэхэд inventory_items + menu_items\nхүснэгтэнд зэрэг бичигдэнэ",
    "Customer тухайн барааг захиалахад GREATEST(qty-delta, 0)\nтомьёогоор хасалт хийгдэнэ",
    "qty ≤ threshold болох агшинд socket илгээж DB-д\nnotification хадгалагдана",
    "Зөвхөн threshold давсан үед мэдэгдэл (давхардал үгүй)",
]
for i, line in enumerate(flow):
    top = Inches(2.45) + i * Inches(0.85)
    num_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.1),
                                    Inches(0.4), Inches(0.4))
    num_circle.line.fill.background()
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = PRIMARY
    add_text(s, Inches(0.7), top + Inches(0.12), Inches(0.4), Inches(0.4),
             str(i + 1), size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.25), top, Inches(6.8), Inches(0.85), line, size=13)

# right column - schema snippet
add_card(s, Inches(8.3), Inches(1.9), Inches(4.45), Inches(4.9))
add_text(s, Inches(8.5), Inches(2.05), Inches(4.1), Inches(0.5),
         "inventory_items", size=15, bold=True, color=PRIMARY,
         font="Consolas")
schema = [
    "id            serial PK",
    "name          text",
    "type          text",
    "quantity      integer",
    "threshold     integer",
    "image_url     text?",
    "created_at    timestamp",
    "updated_at    timestamp",
]
for i, row in enumerate(schema):
    add_text(s, Inches(8.5), Inches(2.55) + i * Inches(0.38), Inches(4.1),
             Inches(0.35), row, size=12, color=WHITE, font="Consolas")
add_text(s, Inches(8.5), Inches(5.8), Inches(4.1), Inches(0.4),
         "menu_items.inventory_item_id → FK",
         size=11, color=MUTED, font="Consolas")
footer(s, 9, 12)

# ─── Slide 10: Notifications ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Бодит цагийн мэдэгдлийн систем")

add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Event → DB → Socket → UI", size=16, bold=True, color=PRIMARY)

# Flow diagram (boxes + arrows)
boxes = [
    ("Event", "Шинэ захиалга\nСэтгэгдэл\nНөөц дуусах"),
    ("createNotification()", "DB-д бичиж,\nSocket.IO-р илгээнэ"),
    ("notifications\nхүснэгт", "type · title · message\ndata · read"),
    ("UI хонх", "Улаан бадж +\nунших / цэвэрлэх"),
]
bw = Inches(2.8)
bh = Inches(2.0)
gap = Inches(0.35)
total = 4 * bw + 3 * gap
start = (SW - total) / 2
top = Inches(2.75)
for i, (t, d) in enumerate(boxes):
    left = start + i * (bw + gap)
    add_card(s, left, top, bw, bh)
    add_text(s, left + Inches(0.2), top + Inches(0.3), bw - Inches(0.4),
             Inches(0.6), t, size=16, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.2), top + Inches(1.0), bw - Inches(0.4),
             Inches(0.9), d, size=12, align=PP_ALIGN.CENTER)
    if i < 3:
        ax = left + bw
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Emu(10000),
                                 top + Inches(0.85), gap - Emu(20000),
                                 Inches(0.3))
        arr.line.fill.background()
        arr.fill.solid()
        arr.fill.fore_color.rgb = PRIMARY

add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
         "Эвент төрөл", size=14, bold=True, color=PRIMARY)
types = [
    ("order_new", "Шинэ захиалга үүсэхэд"),
    ("inventory_low", "Нөөц threshold давахад"),
    ("review_new", "Зочин сэтгэгдэл илгээхэд"),
]
for i, (t, d) in enumerate(types):
    top2 = Inches(5.8) + i * Inches(0.4)
    add_text(s, Inches(0.9), top2, Inches(3), Inches(0.35), t, size=12,
             bold=True, color=ACCENT, font="Consolas")
    add_text(s, Inches(3.8), top2, Inches(8), Inches(0.35), "→ " + d,
             size=12, color=WHITE)
footer(s, 10, 12)

# ─── Slide 11: Results ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "Гарсан үр дүн")

# stat cards
stats = [
    ("9", "үндсэн модуль"),
    ("14+", "REST endpoint"),
    ("4", "ажилтны роль"),
    ("∞", "зэрэг зочин / ширээ"),
]
cw = Inches(2.9)
gap = Inches(0.2)
left0 = Inches(0.6)
top = Inches(2.1)
for i, (n, label) in enumerate(stats):
    left = left0 + i * (cw + gap)
    add_card(s, left, top, cw, Inches(1.8))
    add_text(s, left, top + Inches(0.25), cw, Inches(0.9),
             n, size=44, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, left, top + Inches(1.2), cw, Inches(0.4), label,
             size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "Хэрэгжсэн шинэ боломжууд", size=16, bold=True, color=PRIMARY)
impl = [
    "Олон зочин нэг идэвхтэй QR дээр зэрэг захиалга өгөх",
    "Гадна талаас орж ирсэн бараа (ундаа, дарс) автомат цэсэнд нэгдэх",
    "Нөөц threshold давахад мэдэгдэх, 0 болсныг ялгах",
    "Бодит цагийн хонхны мэдэгдэл, түүх, унших/цэвэрлэх",
    "Ширээ тус бүрийн идэвхтэй захиалгын нийт дүнг нэгтгэн харуулах",
    "Топ хоолыг огноогоор шүүж бүх төрлийг харуулах тайлан",
]
for i, line in enumerate(impl):
    top_i = Inches(4.8) + i * Inches(0.38)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), top_i + Inches(0.12),
                             Inches(0.12), Inches(0.12))
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    add_text(s, Inches(1.0), top_i, Inches(11.5), Inches(0.4), line, size=13)
footer(s, 11, 12)

# ─── Slide 12: Thanks ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SH)
stripe.line.fill.background()
stripe.fill.solid()
stripe.fill.fore_color.rgb = PRIMARY

add_text(s, Inches(1.0), Inches(2.2), Inches(12), Inches(1.2),
         "Анхаарал хандуулсанд\nБаярлалаа", size=54, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(4.4), Inches(12), Inches(0.5),
         "Асуулт хариулт", size=22, color=PRIMARY, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(5.1), Inches(12), Inches(0.5),
         "GitHub: github.com/Ochko0310/QR_Restaurant", size=14,
         color=MUTED, align=PP_ALIGN.CENTER)

out = "QR_Restaurant_Presentation.pptx"
prs.save(out)
print("saved", out)
