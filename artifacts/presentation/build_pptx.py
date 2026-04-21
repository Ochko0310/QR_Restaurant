# -*- coding: utf-8 -*-
"""
Рестораны ширээ захиалга, цэс удирдлагын мобайл систем
Дипломын хамгаалалтын илтгэлийн слайд (PowerPoint, 16:9).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ======================= Өнгөний схем =======================
PRIMARY = RGBColor(0x00, 0x46, 0x7F)     # MUSTBlue
ACCENT  = RGBColor(0x46, 0x82, 0xB4)     # Steel blue
OLIVE   = RGBColor(0x55, 0x6B, 0x2F)     # Olive
GOLD    = RGBColor(0xD4, 0xA0, 0x17)     # Accent gold
DARK    = RGBColor(0x1E, 0x1E, 0x1E)
GRAY    = RGBColor(0x60, 0x60, 0x60)
LIGHT   = RGBColor(0xF2, 0xF6, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
RED     = RGBColor(0xC6, 0x28, 0x28)

FONT = "Calibri"  # Monghol-compatible on Windows; fallback widely available

# ======================= Үндсэн тохируулга =======================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ======================= Туслах функцууд =======================
def add_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s

def add_round_rect(slide, x, y, w, h, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    # Дээд туг
    add_rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11.5), Inches(0.7),
             title, size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Баруун буланд логошинх
    add_text(slide, Inches(10.8), Inches(0.18), Inches(2.3), Inches(0.35),
             "Өмнөговь ТДС", size=11, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.8), Inches(0.5), Inches(2.3), Inches(0.3),
             "Бакалаврын дипломын ажил", size=9, italic=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
                 subtitle, size=14, italic=True, color=ACCENT)

def footer(slide, page_num, total):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), PRIMARY)
    add_text(slide, Inches(0.4), SH - Inches(0.35), Inches(8), Inches(0.35),
             "Рестораны ширээ захиалга, цэс удирдлагын мобайл систем",
             size=9, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2.2), SH - Inches(0.35), Inches(2.0), Inches(0.35),
             f"{page_num} / {total}", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def bullet_list(slide, x, y, w, h, items, *, size=16, bold_lead=True, color=DARK,
                lead_color=PRIMARY, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Гол цэг
        r0 = p.add_run()
        r0.text = "▸ "
        r0.font.name = FONT; r0.font.size = Pt(size)
        r0.font.bold = True; r0.font.color.rgb = lead_color
        # Тайлбар
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.name = FONT; r1.font.size = Pt(size)
            r1.font.bold = bold_lead; r1.font.color.rgb = lead_color
            r2 = p.add_run()
            r2.text = " — " + rest
            r2.font.name = FONT; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r1 = p.add_run()
            r1.text = str(item)
            r1.font.name = FONT; r1.font.size = Pt(size)
            r1.font.color.rgb = color
    return tb

def card(slide, x, y, w, h, title, body, *, title_color=PRIMARY, bg=LIGHT):
    add_round_rect(slide, x, y, w, h, bg, line_color=ACCENT)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.45),
             title, size=15, bold=True, color=title_color, anchor=MSO_ANCHOR.TOP)
    add_text(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), h - Inches(0.65),
             body, size=12, color=DARK)

def stat_card(slide, x, y, w, h, value, label, *, val_color=PRIMARY, bg=WHITE):
    add_round_rect(slide, x, y, w, h, bg, line_color=ACCENT)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.9),
             value, size=34, bold=True, color=val_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(1.05), w, h - Inches(1.1),
             label, size=11, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

def make_table(slide, x, y, w, h, headers, rows, *, header_color=PRIMARY,
               zebra=True, header_size=12, body_size=11, col_widths=None):
    cols = len(headers); rws = len(rows) + 1
    tbl_shape = slide.shapes.add_table(rws, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    # Header
    for c, h_ in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(54000); tf.margin_right = Emu(54000)
        tf.margin_top = Emu(36000); tf.margin_bottom = Emu(36000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for r in list(p.runs): r._r.getparent().remove(r._r)
        run = p.add_run(); run.text = h_
        run.font.name = FONT; run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = WHITE
    # Body
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            if zebra and ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = Emu(54000); tf.margin_right = Emu(54000)
            tf.margin_top = Emu(27000); tf.margin_bottom = Emu(27000)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            for r in list(p.runs): r._r.getparent().remove(r._r)
            run = p.add_run(); run.text = str(val)
            run.font.name = FONT; run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
    return tbl

def add_chart(slide, chart_type, x, y, w, h, categories, series, *,
              title=None, legend=True, palette=None):
    """Графикыг слайд дээр нэмэх туслах функц."""
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    gframe = slide.shapes.add_chart(chart_type, x, y, w, h, data)
    chart = gframe.chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for p in chart.chart_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = PRIMARY
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        for p in chart.legend.font._element.getparent() if False else []:
            pass
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = FONT
    # Цуврал бүрийн өнгө
    if palette:
        for i, plot in enumerate(chart.plots):
            for j, ser in enumerate(plot.series):
                idx = (i * len(plot.series) + j) % len(palette)
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = palette[idx]
                ser.format.line.color.rgb = palette[idx]
    return chart

def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, weight=2.5):
    cn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # STRAIGHT
    cn.line.color.rgb = color
    cn.line.width = Pt(weight)
    # Arrow head
    ln = cn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med"); tail.set("len", "med")
    return cn

# ======================= СЛАЙДУУД =======================
SLIDES = []

# -------- Slide 1: Нүүр --------
def slide_cover():
    s = add_slide()
    set_bg(s, PRIMARY)
    # Алтан зурвас
    add_rect(s, 0, Inches(0.0), SW, Inches(0.15), GOLD)
    add_rect(s, 0, SH - Inches(0.15), SW, Inches(0.15), GOLD)
    # Сургууль
    add_text(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5),
             "ӨМНӨГОВЬ АЙМАГ ДАХЬ ТЕХНОЛОГИЙН ДЭЭД СУРГУУЛЬ",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
             "Мэдээллийн технологийн тэнхим",
             size=13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Гол гарчиг
    add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.3),
             "РЕСТОРАНЫ ШИРЭЭ ЗАХИАЛГА,\nЦЭС УДИРДЛАГЫН МОБАЙЛ СИСТЕМ",
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Мөр
    add_rect(s, Inches(4), Inches(3.9), Inches(5.3), Inches(0.04), GOLD)
    # Дэд гарчиг
    add_text(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5),
             "React Native · Node.js · PostgreSQL · Socket.IO дээр суурилсан",
             size=16, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # Оюутан / удирдагч
    add_text(s, Inches(1), Inches(5.3), Inches(5.5), Inches(0.4),
             "Дипломант:", size=13, bold=True, color=GOLD)
    add_text(s, Inches(1), Inches(5.7), Inches(5.5), Inches(0.4),
             "Эрдэнэтөгсийн Очбадрах", size=16, bold=True, color=WHITE)
    add_text(s, Inches(1), Inches(6.1), Inches(5.5), Inches(0.4),
             "IT-XX бүлэг", size=12, color=LIGHT)
    add_text(s, Inches(7), Inches(5.3), Inches(5.5), Inches(0.4),
             "Удирдагч багш:", size=13, bold=True, color=GOLD)
    add_text(s, Inches(7), Inches(5.7), Inches(5.5), Inches(0.4),
             "Б. Сэргэлэнбаяр", size=16, bold=True, color=WHITE)
    add_text(s, Inches(7), Inches(6.1), Inches(5.5), Inches(0.4),
             "Ахлах багш", size=12, color=LIGHT)
    # Огноо
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "Даланзадгад · 2026", size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

SLIDES.append(slide_cover)

# -------- Slide 2: Илтгэлийн бүтэц --------
def slide_agenda():
    s = add_slide(); header(s, "Илтгэлийн бүтэц")
    items = [
        "1. Судалгааны сэдвийн үндэслэл, асуудал, зорилго",
        "2. Одоогийн шийдлүүдтэй харьцуулсан судалгаа",
        "3. Функциональ ба функциональ бус шаардлагууд",
        "4. Системийн архитектур, технологи сонголт",
        "5. Хэрэглэгчийн хэрэглээний диаграм (Use Case)",
        "6. Өгөгдлийн сангийн загвар (ER диаграм)",
        "7. Гол бизнес процессуудын дараалал (Sequence)",
        "8. API, WebSocket, аюулгүй байдал",
        "9. Байршуулалтын архитектур (Deployment)",
        "10. Туршилт, үр дүн",
        "11. Дүгнэлт, цаашдын ажил",
    ]
    # Хоёр багана
    left = items[:6]; right = items[6:]
    tbL = slide_column(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.4), left)
    tbR = slide_column(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4), right)

def slide_column(slide, x, y, w, h, items):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(54000); tf.margin_right = Emu(54000)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.5; p.space_after = Pt(8)
        r = p.add_run(); r.text = it
        r.font.name = FONT; r.font.size = Pt(17); r.font.color.rgb = DARK
    return tb

SLIDES.append(slide_agenda)

# -------- Slide 3: Үндэслэл --------
def slide_background():
    s = add_slide(); header(s, "Судалгааны үндэслэл")
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
             "Монголын зоогийн газрын зах зээлийн дижитал шилжилтийн бэлэн байдал",
             size=15, italic=True, color=ACCENT)
    # 4 statistic cards
    stat_card(s, Inches(0.5),  Inches(1.9), Inches(3), Inches(1.9), "4,200+",
              "Монгол дахь нийт зоогийн газар", val_color=PRIMARY)
    stat_card(s, Inches(3.6),  Inches(1.9), Inches(3), Inches(1.9), "78%",
              "Захиалгыг гар аргаар авдаг", val_color=RED)
    stat_card(s, Inches(6.7),  Inches(1.9), Inches(3), Inches(1.9), "62%",
              "Захиалгын дундаж хугацаа хэмнэгдэх", val_color=GREEN)
    stat_card(s, Inches(9.8),  Inches(1.9), Inches(3), Inches(1.9), "~90%",
              "Smartphone-ын нэвтрэлт", val_color=OLIVE)
    # Тайлбар
    bullet_list(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0), [
        ("Олон улсын системүүд",
         "сарын 200–500$+ төлбөртэй, Монгол хэл, QPay/SocialPay-г бүрэн дэмждэггүй."),
        ("Орон нутгийн систем",
         "хүрэлцээгүй, ихэнх нь зөвхөн вэб тулгуурт, мобайл хэрэглэгчид тохирохгүй."),
        ("COVID-19 дараах эрэлт",
         "QR-аар захиалах, контакгүй тооцоо — хэрэглэгчид дадсан."),
        ("Боломж",
         "Хямд, монгол хэлтэй, мобайл тэргүүлэх, бодит цагийн шийдэл маш бага."),
    ], size=14)

SLIDES.append(slide_background)

# -------- Slide 4: Асуудал --------
def slide_problem():
    s = add_slide(); header(s, "Одоогийн шийдлүүдийн асуудал")
    rows = [
        ("1", "Гадаадын систем үнэтэй",        "Toast, Square — сарын 200–500 ам.доллар; жижиг зоогийн газарт хэт өндөр."),
        ("2", "Монгол хэлний дэмжлэг хомс",    "Олон улсын системүүд UI-г монголчилдоггүй, мэдэгдэл англиар."),
        ("3", "Төлбөрийн интеграц дутмаг",     "QPay/SocialPay/Monpay-г холбоогүй; зөвхөн Stripe/PayPal."),
        ("4", "Мобайл тэргүүлэх шийдэл бага",  "Ихэнх нь вэб; таблет POS дээр тулгуурласан, гар утсаар хэцүү."),
        ("5", "Бодит цагийн шинэчлэл сул",     "Захиалга, статус удаан ирдэг; гал тогоо–зааланд мэдээлэл алдагдах."),
        ("6", "Орон нутгийн дэмжлэг байхгүй",  "Суурилуулалт, сургалт, дэмжлэг орон нутагт хүрдэггүй."),
    ]
    make_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.6),
               ["#", "Асуудал", "Тодорхойлолт"], rows,
               col_widths=[1, 4, 10], body_size=12, header_size=13)

SLIDES.append(slide_problem)

# -------- Slide 4b: Системүүдийн үнийн харьцуулалт --------
def slide_price_chart():
    s = add_slide(); header(s, "Сарын үнийн харьцуулалт (ам.доллар)")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "Жижиг зоогийн газарт зориулсан тарифаар (median restaurant, 30 ширээ)",
             size=13, italic=True, color=ACCENT)
    add_chart(
        s, XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(1.7), Inches(8.0), Inches(5.2),
        categories=["Toast POS", "Square", "Lightspeed", "TouchBistro",
                    "Resy", "Монгол систем (энэ)"],
        series=[("Сарын үнэ (USD)", (165, 129, 189, 249, 899, 18))],
        title="Үнийн харьцуулалт",
        legend=False,
        palette=[RED, RED, RED, RED, RED, GREEN],
    )
    # Тайлбар
    bullet_list(s, Inches(8.7), Inches(1.7), Inches(4.3), Inches(5.2), [
        ("Санхүү", "Барууны POS сар 130–900$; Монголын SME-д дийлдэхгүй."),
        ("Манай систем", "~18 USD (hosting + QPay комисс)."),
        ("Хэмнэлт", "90–98% хямд — ROI 4-6 сард."),
        ("Дэлхийн стандарт", "Үндсэн функцууд бүгд орсон."),
    ], size=12)

SLIDES.append(slide_price_chart)

# -------- Slide 5: Зорилго, зорилт --------
def slide_goal():
    s = add_slide(); header(s, "Зорилго ба зорилтууд")
    # Зорилгын хайрцаг
    add_round_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.4),
                   PRIMARY, line_color=PRIMARY)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
             "ЗОРИЛГО", size=13, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0),
             ("Монголын зоогийн газрын онцлогт тохирсон, QR-аар ширээ захиалах,\n"
              "цэс удирдах, бодит цагийн захиалгын зохицуулалттай, QPay-тэй\n"
              "интеграц хийсэн mobile-first систем бүтээж, нэвтрүүлэх."),
             size=14, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 6 зорилт
    goals = [
        ("1. Үнэлгээ", "Одоогийн системүүдийн харьцуулсан шинжилгээ, оролцогчдын судалгаа."),
        ("2. Шаардлага", "Функциональ (12), функциональ бус (ISO/IEC 25010, 5 ангилал)."),
        ("3. Архитектур", "Layered + monorepo + бүрэлдэхүүн хэсгүүдийн ялгаралт."),
        ("4. Хөгжүүлэлт", "Mobile (React Native), API (Express), Socket.IO, PostgreSQL."),
        ("5. Туршилт", "Функциональ, интеграц, ачаалал (k6), хэрэглэхүйн (SUS)."),
        ("6. Үнэлгээ", "Хугацаа, алдаа, сэтгэл ханамжийн бодит үзүүлэлт.")
    ]
    y = Inches(2.9); cw = Inches(3.95); ch = Inches(1.9); gap = Inches(0.15)
    for i, (t, d) in enumerate(goals):
        col = i % 3; row = i // 3
        x = Inches(0.5) + col * (cw + gap)
        yy = y + row * (ch + gap)
        card(s, x, yy, cw, ch, t, d)

SLIDES.append(slide_goal)

# -------- Slide 6: Шинэлэг тал, хязгаарлалт --------
def slide_novelty():
    s = add_slide(); header(s, "Шинэлэг тал ба судалгааны хязгаар")
    # Зүүн: шинэлэг тал
    add_round_rect(s, Inches(0.5), Inches(1.2), Inches(6.15), Inches(5.8),
                   LIGHT, line_color=ACCENT)
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.45),
             "Шинэлэг тал", size=17, bold=True, color=PRIMARY)
    bullet_list(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(5.1), [
        ("Монгол хэл", "Бүх UI, имэйл, нэхэмжлэх монгол."),
        ("QPay интеграц", "Callback-аар автоматаар тооцоо хаагдана."),
        ("Mobile-first", "React Native нэг codebase — iOS + Android."),
        ("Socket.IO бодит цагт", "Захиалга, статусын шинэчлэл 1 сек-д."),
        ("Хэрэглэхүйн судалгаа", "SUS шкалаар эцсийн хэрэглэгчид үнэлгээ."),
        ("Нээлттэй архитектур", "Layered + RBAC + Drizzle ORM → засварлах хялбар."),
    ], size=13)
    # Баруун: хязгаарлалт
    add_round_rect(s, Inches(6.85), Inches(1.2), Inches(6.15), Inches(5.8),
                   LIGHT, line_color=RED)
    add_text(s, Inches(7.05), Inches(1.3), Inches(5.8), Inches(0.45),
             "Хязгаарлалт", size=17, bold=True, color=RED)
    bullet_list(s, Inches(7.05), Inches(1.8), Inches(5.8), Inches(5.1), [
        ("Цар хүрээ", "1 зоогийн газар дээр прототип, олон салбар дэмжлэг II шатанд."),
        ("Нөөц удирдлага", "Түүхий эдийн анхан шатны хяналт, нарийн нягтлан хөтлөлтгүй."),
        ("AR цэс", "Зөвхөн зураг/видео; 3D/AR дараагийн хувилбарт."),
        ("Хүргэлт", "Зөвхөн зоогийн газар доторх захиалга, delivery II шатанд."),
        ("Сэтгэгдлийн ML", "Рэйтинг бичлэг; автомат sentiment-ийг үгүйсгэсэн."),
        ("Оффлайн горим", "Зөвхөн урьдчилан татсан цэс, захиалга онлайн."),
    ], size=13)

SLIDES.append(slide_novelty)

# -------- Slide 7: Функциональ шаардлага --------
def slide_fr():
    s = add_slide(); header(s, "Функциональ шаардлагууд (FR)")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "MoSCoW эрэмбэтэй 12 үндсэн шаардлага",
             size=13, italic=True, color=ACCENT)
    rows = [
        ("FR-01", "QR-аар цэс нээх",            "M"),
        ("FR-02", "Ширээ урьдчилан захиалах",   "M"),
        ("FR-03", "Онлайн захиалга үүсгэх",     "M"),
        ("FR-04", "Бодит цагийн мэдэгдэл",      "M"),
        ("FR-05", "QPay-ээр төлбөр хийх",       "M"),
        ("FR-06", "Цэс/бараа CRUD удирдлага",   "M"),
        ("FR-07", "Захиалгын төлөв шинэчлэх",   "M"),
        ("FR-08", "JWT + RBAC аутентикация",    "M"),
        ("FR-09", "Ажилтны панель (кассир)",    "S"),
        ("FR-10", "Менежерийн тайлан",          "S"),
        ("FR-11", "Сэтгэгдэл, үнэлгээ",         "S"),
        ("FR-12", "Нөөцийн автомат хасалт",     "C"),
    ]
    make_table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.0),
               ["Код", "Шаардлага", "MoSCoW"], rows,
               col_widths=[2, 9, 2], header_size=13, body_size=13)
    # Тэмдэглэгээ
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "M = Must (заавал), S = Should (хэрэгтэй), C = Could (хүсэмжит)",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_fr)

# -------- Slide 8: Функциональ бус шаардлага --------
def slide_nfr():
    s = add_slide(); header(s, "Функциональ бус шаардлага (ISO/IEC 25010)")
    cats = [
        ("Гүйцэтгэл",     "API p95 < 300мс · WebSocket < 100мс · 200 зэрэг VU · DB QPS ≥ 150"),
        ("Найдвартай\nбайдал", "Uptime ≥ 99.5% · MTBF > 720ц · RPO 15мин · RTO 2ц · Auto-restart"),
        ("Аюулгүй\nбайдал",   "TLS 1.3 · bcrypt(12) · JWT (15мин) · RBAC · OWASP Top-10 mitigated"),
        ("Хэрэглэхүй",   "SUS > 75 · WCAG 2.1 AA · Монгол/англи · Touch target ≥ 44px"),
        ("Засварлах",    "TypeScript · ESLint · Test coverage > 70% · OpenAPI docs"),
    ]
    y = Inches(1.3); cw = Inches(12.3); ch = Inches(1.05); gap = Inches(0.12)
    colors = [PRIMARY, OLIVE, RED, ACCENT, GOLD]
    for i, (t, d) in enumerate(cats):
        add_round_rect(s, Inches(0.5), y, Inches(2.5), ch, colors[i])
        add_text(s, Inches(0.5), y, Inches(2.5), ch, t,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round_rect(s, Inches(3.1), y, Inches(9.7), ch, LIGHT, line_color=ACCENT)
        add_text(s, Inches(3.25), y, Inches(9.55), ch, d,
                 size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + ch + gap

SLIDES.append(slide_nfr)

# -------- Slide 9: Use Case --------
def slide_usecase():
    s = add_slide(); header(s, "Хэрэглэгчийн хэрэглээний диаграм (Use Case)")
    # 4 actor
    actors = [
        (Inches(0.3),  Inches(3.3), "Хэрэглэгч"),
        (Inches(0.3),  Inches(5.5), "Зочин\n(QR)"),
        (Inches(11.5), Inches(2.6), "Ажилтан\n(кассир/менежер)"),
        (Inches(11.5), Inches(5.5), "Гал тогоо"),
    ]
    for (x, y, t) in actors:
        # stickman rectangle
        add_round_rect(s, x, y, Inches(1.5), Inches(1.6), ACCENT, line_color=ACCENT)
        add_text(s, x, y, Inches(1.5), Inches(1.6), t,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # System ellipse container
    add_round_rect(s, Inches(2.2), Inches(1.4), Inches(9.0), Inches(5.5),
                   LIGHT, line_color=PRIMARY)
    add_text(s, Inches(2.2), Inches(1.45), Inches(9.0), Inches(0.35),
             "Системийн хил", size=11, bold=True, italic=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    # Use cases (эллипс шиг round rect)
    ucs = [
        ("ЮК-01 QR уншуулах",         Inches(2.45), Inches(1.95)),
        ("ЮК-02 Цэс үзэх",            Inches(4.85), Inches(1.95)),
        ("ЮК-03 Захиалга өгөх",       Inches(7.25), Inches(1.95)),
        ("ЮК-04 Ширээ сонгох",        Inches(2.45), Inches(2.75)),
        ("ЮК-05 QPay төлбөр",         Inches(4.85), Inches(2.75)),
        ("ЮК-06 Бүртгэл/нэвтрэлт",    Inches(7.25), Inches(2.75)),
        ("ЮК-07 Захиалга баталгаажуулах", Inches(2.45), Inches(3.55)),
        ("ЮК-08 Цэс удирдах",         Inches(4.85), Inches(3.55)),
        ("ЮК-09 Тайлан харах",        Inches(7.25), Inches(3.55)),
        ("ЮК-10 Нөөц удирдах",        Inches(2.45), Inches(4.35)),
        ("ЮК-11 Хоол бэлтгэх",        Inches(4.85), Inches(4.35)),
        ("ЮК-12 Төлөв шинэчлэх",      Inches(7.25), Inches(4.35)),
    ]
    for t, x, y in ucs:
        add_round_rect(s, x, y, Inches(2.3), Inches(0.65), WHITE, line_color=PRIMARY)
        add_text(s, x, y, Inches(2.3), Inches(0.65), t,
                 size=10, color=DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Тайлбар
    add_text(s, Inches(2.2), Inches(5.2), Inches(9.0), Inches(1.5),
             ("include: «QR уншуулах» → «Цэс үзэх»; «Захиалга өгөх» → «QPay төлбөр»\n"
              "extend: «Захиалга өгөх» ← «Ширээ сонгох» (зоогийн газар дотор)\n"
              "Нийт 4 role × 12 use case = 48 эрхийн холбоос (RBAC матриц)."),
             size=11, italic=True, color=GRAY, align=PP_ALIGN.LEFT)

SLIDES.append(slide_usecase)

# -------- Slide 10: RBAC матриц --------
def slide_rbac():
    s = add_slide(); header(s, "Хандалтын эрхийн матриц (RBAC)")
    headers = ["Үүрэг", "Цэс", "Захиалга", "Ширээ", "Төлбөр", "Хэрэглэгч", "Тайлан", "Нөөц"]
    rows = [
        ("Зочин",             "R",   "C",   "R",   "C",   "—",   "—",   "—"),
        ("Хэрэглэгч",         "R",   "CRU", "CR",  "C",   "RU*", "R*",  "—"),
        ("Ажилтан",           "R",   "CRU", "CRU", "CRU", "R",   "R",   "RU"),
        ("Менежер / Admin",   "CRUD","CRUD","CRUD","CRUD","CRUD","CRUD","CRUD"),
        ("Гал тогоо",         "R",   "RU",  "—",   "—",   "—",   "—",   "RU"),
    ]
    make_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5),
               headers, rows, header_size=12, body_size=12,
               col_widths=[3, 1.6, 1.6, 1.6, 1.6, 1.8, 1.6, 1.6])
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1),
             ("C = Create · R = Read · U = Update · D = Delete · —  нэвтрэх эрхгүй\n"
              "* Өөрийн бүртгэл/захиалгыг л уншиж/засварлах эрхтэй."),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_rbac)

# -------- Slide 11: Системийн архитектур --------
def slide_arch():
    s = add_slide(); header(s, "Системийн архитектур (Layered)")
    layers = [
        ("Presentation Layer",
         "React Native (Expo) · iOS/Android · React Navigation · Zustand",
         PRIMARY),
        ("Application Layer",
         "Node.js · Express · Socket.IO · JWT middleware · RBAC",
         ACCENT),
        ("Business Logic",
         "Захиалга, менежмент, тооцоо, нөөц хасалт, гүйлгээ",
         OLIVE),
        ("Data Access",
         "Drizzle ORM · Migration · ACID гүйлгээ",
         GOLD),
        ("Data Layer",
         "PostgreSQL 16 · Redis (cache/session) · S3 (зураг)",
         RED),
    ]
    y = Inches(1.3); ch = Inches(1.0); gap = Inches(0.08)
    for t, d, c in layers:
        add_rect(s, Inches(1.5), y, Inches(10.3), ch, c)
        add_text(s, Inches(1.5), y, Inches(3.2), ch, t,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.8), y, Inches(6.9), ch, d,
                 size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        y = y + ch + gap
    # Гадаад систем
    add_round_rect(s, Inches(0.3), Inches(1.3), Inches(1.0), Inches(5.4), LIGHT, line_color=PRIMARY)
    add_text(s, Inches(0.3), Inches(1.3), Inches(1.0), Inches(5.4),
             "Гадны\nсистем\n—\nQPay\nSMTP\nFCM",
             size=10, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

SLIDES.append(slide_arch)

# -------- Slide 12: Технологи сонголт --------
def slide_tech():
    s = add_slide(); header(s, "Технологи сонголт ба үндэслэл")
    rows = [
        ("Давхарга", "Технологи", "Яагаад сонгосон"),
    ]
    # using main make_table pattern
    data = [
        ("Мобайл", "React Native + Expo",
         "iOS + Android нэг codebase · OTA шинэчлэл · дэлхийн оюутан сурах нөөц өргөн"),
        ("Backend", "Node.js + Express",
         "Non-blocking I/O · JSON-тай байгалиасаа · npm ecosystem · хурдан прототип"),
        ("Realtime", "Socket.IO",
         "WebSocket + fallback · room/namespace · Express-тэй нэгдмэл"),
        ("ORM", "Drizzle ORM",
         "TypeScript-first · SQL-с ойр · migration control · type-safe query"),
        ("Database", "PostgreSQL 16",
         "ACID · JSONB · partial index · Mongol hiragana-тэй collation"),
        ("Cache", "Redis",
         "Session · rate-limit · pub/sub (олон instance дахь Socket.IO)"),
        ("Payment", "QPay REST API",
         "Монголын зах зээл · callback гүйлгээ · QR-н invoice"),
        ("DevOps", "Docker + Nginx",
         "Orchestration хөнгөн · TLS termination · 1 VPS-д багтдаг"),
    ]
    make_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
               ["Давхарга", "Технологи", "Үндэслэл"], data,
               col_widths=[2.5, 3.5, 8], header_size=13, body_size=12)

SLIDES.append(slide_tech)

# -------- Slide 13: Monorepo / багцын бүтэц --------
def slide_monorepo():
    s = add_slide(); header(s, "Monorepo бүтэц (pnpm workspaces)")
    pkgs = [
        ("apps/mobile",         "React Native аппликейшн (хэрэглэгч, ажилтан)"),
        ("apps/admin-web",      "React + Vite менежерийн вэб панель"),
        ("services/api-server", "Express REST + Socket.IO сервер"),
        ("services/worker",     "BullMQ queue: имэйл, мэдэгдэл, тайлан"),
        ("packages/shared",     "TypeScript interface, Zod schema, types"),
        ("packages/ui",         "React Native нэгдсэн UI компонент"),
        ("packages/db",         "Drizzle ORM schema + migration скрипт"),
        ("infra/docker",        "Dockerfile, docker-compose, Nginx config"),
    ]
    make_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0),
               ["Багц (package)", "Зориулалт"], pkgs,
               col_widths=[4, 9], header_size=13, body_size=13)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
             "pnpm хурдан суулгалт · workspace-аар хамаарал хуваалцах · тусгай CI cache.",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_monorepo)

# -------- Slide 14: ER диаграм --------
def slide_er():
    s = add_slide(); header(s, "Өгөгдлийн сангийн загвар (ER)")
    ents = [
        ("users",          Inches(0.4),  Inches(1.3), PRIMARY),
        ("restaurants",    Inches(4.8),  Inches(1.3), PRIMARY),
        ("tables",         Inches(9.2),  Inches(1.3), PRIMARY),
        ("menu_items",     Inches(0.4),  Inches(3.0), ACCENT),
        ("categories",     Inches(4.8),  Inches(3.0), ACCENT),
        ("orders",         Inches(9.2),  Inches(3.0), ACCENT),
        ("order_items",    Inches(0.4),  Inches(4.7), OLIVE),
        ("inventory_items",Inches(4.8),  Inches(4.7), OLIVE),
        ("table_sessions", Inches(9.2),  Inches(4.7), OLIVE),
    ]
    for (t, x, y, c) in ents:
        add_round_rect(s, x, y, Inches(3.6), Inches(1.3), c)
        add_text(s, x, y + Inches(0.05), Inches(3.6), Inches(0.45),
                 t, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Гол талбар
        fields = {
            "users": "id · role · phone · email\npassword_hash · created_at",
            "restaurants": "id · name · address · phone\nopen_hours · owner_id",
            "tables": "id · restaurant_id · number\ncapacity · qr_code · status",
            "menu_items": "id · category_id · name · price\nimage_url · is_available",
            "categories": "id · restaurant_id\nname · display_order",
            "orders": "id · user_id · table_id · status\ntotal · payment_status",
            "order_items": "id · order_id · menu_item_id\nquantity · price",
            "inventory_items": "id · restaurant_id · name\nquantity · unit · alert",
            "table_sessions": "id · table_id · start/end_time\nguest_count · total",
        }
        add_text(s, x + Inches(0.1), y + Inches(0.5), Inches(3.4), Inches(0.8),
                 fields[t], size=9, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # Хамаарлын сум
    rels = [
        (Inches(2.2), Inches(1.95), Inches(4.8),  Inches(1.95)),   # restaurants-users
        (Inches(6.6), Inches(1.95), Inches(9.2),  Inches(1.95)),   # restaurants-tables
        (Inches(6.6), Inches(2.6),  Inches(6.6),  Inches(3.0)),    # restaurants-categories
        (Inches(2.2), Inches(3.65), Inches(4.8),  Inches(3.65)),   # menu_items-categories
        (Inches(6.6), Inches(3.65), Inches(9.2),  Inches(3.65)),   # categories-orders
        (Inches(2.2), Inches(5.35), Inches(9.2),  Inches(3.65)),   # order_items-orders
        (Inches(11.0),Inches(2.6),  Inches(11.0), Inches(4.7)),    # tables-table_sessions
    ]
    for (x1,y1,x2,y2) in rels:
        add_arrow(s, x1, y1, x2, y2, color=GRAY, weight=1.5)

SLIDES.append(slide_er)

# -------- Slide 15: Гол хүснэгтүүд --------
def slide_schema():
    s = add_slide(); header(s, "Гол хүснэгтийн бүдүүвч")
    # 3 columns
    defs = [
        ("orders", PRIMARY, [
            "id UUID PK",
            "user_id → users.id",
            "table_id → tables.id",
            "status enum(",
            "  pending|preparing|",
            "  serving|completed|",
            "  cancelled)",
            "total NUMERIC(10,2)",
            "payment_status enum",
            "created_at TIMESTAMP",
        ]),
        ("table_sessions", ACCENT, [
            "id UUID PK",
            "table_id → tables.id",
            "start_time TIMESTAMP",
            "end_time TIMESTAMP",
            "guest_count INT",
            "total NUMERIC(10,2)",
            "status enum(open|closed)",
            "UNIQUE(table_id, end",
            "IS NULL) — нэг нээлттэй"
        ]),
        ("inventory_items", OLIVE, [
            "id UUID PK",
            "restaurant_id FK",
            "name VARCHAR(120)",
            "quantity NUMERIC",
            "unit VARCHAR(20)",
            "min_quantity NUMERIC",
            "updated_at TIMESTAMP",
            "CHECK(quantity >= 0)",
            "Auto-deduct by trigger",
        ]),
    ]
    for i, (name, col, rows) in enumerate(defs):
        x = Inches(0.5) + i * Inches(4.3)
        add_round_rect(s, x, Inches(1.3), Inches(4.1), Inches(5.6),
                       LIGHT, line_color=col)
        add_rect(s, x, Inches(1.3), Inches(4.1), Inches(0.55), col)
        add_text(s, x, Inches(1.3), Inches(4.1), Inches(0.55),
                 name, size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb = s.shapes.add_textbox(x + Inches(0.2), Inches(1.95),
                                  Inches(3.7), Inches(5.0))
        tf = tb.text_frame; tf.word_wrap = True
        for j, row in enumerate(rows):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = row
            r.font.name = "Consolas"; r.font.size = Pt(12); r.font.color.rgb = DARK

SLIDES.append(slide_schema)

# -------- Slide 16: Sequence – Захиалга --------
def slide_sequence_order():
    s = add_slide(); header(s, "Захиалгын дараалал (Sequence)")
    actors = [("Хэрэглэгч", Inches(0.6)),
              ("Мобайл апп", Inches(2.9)),
              ("API сервер", Inches(5.2)),
              ("Socket.IO", Inches(7.6)),
              ("PostgreSQL", Inches(10.0))]
    # lifelines
    for name, x in actors:
        add_rect(s, x, Inches(1.3), Inches(2.1), Inches(0.5), PRIMARY)
        add_text(s, x, Inches(1.3), Inches(2.1), Inches(0.5), name,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # dashed lifeline
        cn = s.shapes.add_connector(1, x + Inches(1.05), Inches(1.8),
                                     x + Inches(1.05), Inches(6.8))
        cn.line.color.rgb = GRAY; cn.line.width = Pt(1)
        ln = cn.line._get_or_add_ln()
        dash = ln.find(qn("a:prstDash"))
        if dash is None:
            dash = etree.SubElement(ln, qn("a:prstDash"))
        dash.set("val", "dash")
    # messages
    msgs = [
        ("1: QR уншуулах",         Inches(0.6)+Inches(1.05), Inches(2.0), Inches(2.9)+Inches(1.05), Inches(2.0)),
        ("2: POST /orders",        Inches(2.9)+Inches(1.05), Inches(2.5), Inches(5.2)+Inches(1.05), Inches(2.5)),
        ("3: validate + auth",     Inches(5.2)+Inches(1.05), Inches(3.0), Inches(5.2)+Inches(1.85), Inches(3.0)),
        ("4: INSERT orders",       Inches(5.2)+Inches(1.05), Inches(3.5), Inches(10.0)+Inches(1.05),Inches(3.5)),
        ("5: return id",           Inches(10.0)+Inches(1.05),Inches(4.0), Inches(5.2)+Inches(1.05), Inches(4.0)),
        ("6: emit «order:new»",    Inches(5.2)+Inches(1.05), Inches(4.5), Inches(7.6)+Inches(1.05), Inches(4.5)),
        ("7: broadcast → staff",   Inches(7.6)+Inches(1.05), Inches(5.0), Inches(2.9)+Inches(1.05), Inches(5.0)),
        ("8: 201 Created JSON",    Inches(5.2)+Inches(1.05), Inches(5.5), Inches(2.9)+Inches(1.05), Inches(5.5)),
        ("9: UI update",           Inches(2.9)+Inches(1.05), Inches(6.0), Inches(0.6)+Inches(1.05), Inches(6.0)),
    ]
    for (t, x1, y1, x2, y2) in msgs:
        add_arrow(s, x1, y1, x2, y2, color=ACCENT, weight=1.5)
        mid_x = min(x1, x2) + Emu(40000)
        add_text(s, mid_x, y1 - Inches(0.30), Inches(4), Inches(0.25),
                 t, size=10, bold=True, color=PRIMARY)

SLIDES.append(slide_sequence_order)

# -------- Slide 17: Socket.IO эвэнтүүд --------
def slide_socket():
    s = add_slide(); header(s, "Socket.IO эвэнт ба бодит цаг")
    rows = [
        ("order:new",        "Клиент → Сервер", "Шинэ захиалга үүсэхэд"),
        ("order:update",     "Сервер → Клиент", "Захиалгын төлөв өөрчлөгдөхөд"),
        ("order:cancelled",  "Сервер → Клиент", "Захиалгыг цуцлахад"),
        ("table:occupied",   "Сервер → Клиент", "Ширээ эзлэгдсэн/чөлөөлөгдсөн"),
        ("kitchen:ready",    "Клиент → Сервер", "Хоол бэлэн болсон мэдэгдэл"),
        ("payment:paid",     "Сервер → Клиент", "QPay callback ирсний дараа"),
        ("menu:updated",     "Сервер → Клиент", "Цэс өөрчлөлт (admin хийх үед)"),
        ("inventory:low",    "Сервер → Клиент", "Нөөцийн анхааруулга (min_qty)"),
        ("staff:notify",     "Сервер → Клиент", "Ажилтанд чиглэсэн мэдэгдэл"),
    ]
    make_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.2),
               ["Эвэнт", "Чиглэл", "Тайлбар"], rows,
               col_widths=[3, 3, 7], header_size=13, body_size=12)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6),
             "Room: restaurant:{id} · user:{id} · table:{id} — зөвхөн хамаатай клиент сонсоно.",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_socket)

# -------- Slide 18: API endpoints --------
def slide_api():
    s = add_slide(); header(s, "REST API endpoint-ууд (30)")
    groups = [
        ("Auth (4)",
         "POST /auth/register · /auth/login · /auth/refresh · /auth/logout"),
        ("Users (4)",
         "GET/PUT /users/me · GET /users/:id · PUT /users/:id/role"),
        ("Restaurants (3)",
         "GET /restaurants · GET /restaurants/:id · PUT /restaurants/:id"),
        ("Tables (4)",
         "GET /tables · POST /tables · PUT /tables/:id · DELETE /tables/:id"),
        ("Menu (5)",
         "GET /menu · POST /menu · PUT/DELETE /menu/:id · PATCH /menu/:id/availability"),
        ("Orders (5)",
         "GET /orders · POST /orders · PATCH /orders/:id/status · /cancel · /items"),
        ("Payments (3)",
         "POST /payments · POST /payments/qpay/callback · GET /payments/:id"),
        ("Inventory (2)",
         "GET /inventory · PATCH /inventory/:id — auto-deduct by order"),
    ]
    y = Inches(1.3); gap = Inches(0.12); ch = Inches(0.64)
    for name, d in groups:
        add_round_rect(s, Inches(0.5),  y, Inches(2.5), ch, PRIMARY)
        add_text(s, Inches(0.5), y, Inches(2.5), ch, name,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round_rect(s, Inches(3.1), y, Inches(9.7), ch, LIGHT, line_color=ACCENT)
        add_text(s, Inches(3.25), y, Inches(9.55), ch, d,
                 size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + ch + gap

SLIDES.append(slide_api)

# -------- Slide 19: Аюулгүй байдал --------
def slide_security():
    s = add_slide(); header(s, "Аюулгүй байдлын шийдлүүд")
    items = [
        ("Authentication", "JWT (HS256) — access 15 мин, refresh 7 хоног. Rotation, blacklist Redis-т."),
        ("Authorization",  "RBAC middleware — эрх шалгалт endpoint бүрт. Cross-tenant filter restaurant_id-аар."),
        ("Password",       "bcrypt cost=12 · password policy (8+, 1 том, 1 тоо, 1 тэмдэгт)."),
        ("Transport",      "TLS 1.3 (Nginx + Let's Encrypt) · HSTS header · secure cookies."),
        ("Input",          "Zod schema validation · parameterized SQL (Drizzle) · HTML-escape (XSS)."),
        ("Rate-limit",     "Express-rate-limit Redis store — login 5/min, API 100/min."),
        ("Payment",        "QPay callback signature verify · idempotent payment_id · HMAC."),
        ("Audit",          "audit_logs хүснэгт — role-өөрчлөлт, нөөц засвар, нэвтрэлт."),
    ]
    # 4x2
    cw = Inches(6.1); ch = Inches(1.28); gap_x = Inches(0.15); gap_y = Inches(0.15)
    for i, (t, d) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.5) + col * (cw + gap_x)
        y = Inches(1.3) + row * (ch + gap_y)
        card(s, x, y, cw, ch, t, d, bg=LIGHT, title_color=RED)

SLIDES.append(slide_security)

# -------- Slide 20: Deployment --------
def slide_deploy():
    s = add_slide(); header(s, "Байршуулалтын архитектур (Deployment)")
    # Cloud хайрцаг
    add_round_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7),
                   LIGHT, line_color=PRIMARY)
    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
             "VPS (Hetzner / DigitalOcean) · Ubuntu 22.04", size=12, bold=True, italic=True,
             color=PRIMARY, align=PP_ALIGN.CENTER)
    # Nginx
    add_round_rect(s, Inches(1.0), Inches(1.9), Inches(2.5), Inches(1.0), PRIMARY)
    add_text(s, Inches(1.0), Inches(1.9), Inches(2.5), Inches(1.0),
             "Nginx\nTLS · reverse proxy", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Docker stack
    add_round_rect(s, Inches(4.5), Inches(1.9), Inches(7.8), Inches(1.0), ACCENT)
    add_text(s, Inches(4.5), Inches(1.9), Inches(7.8), Inches(1.0),
             "Docker Compose · api (x2) · worker · redis",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # API containers
    add_round_rect(s, Inches(4.5), Inches(3.2), Inches(3.6), Inches(1.0), OLIVE)
    add_text(s, Inches(4.5), Inches(3.2), Inches(3.6), Inches(1.0),
             "api-server  (Node.js)\nPort 4000 · stateless",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round_rect(s, Inches(8.6), Inches(3.2), Inches(3.7), Inches(1.0), OLIVE)
    add_text(s, Inches(8.6), Inches(3.2), Inches(3.7), Inches(1.0),
             "worker (BullMQ)\nemail · report · FCM",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # DB primary/replica
    add_round_rect(s, Inches(1.0), Inches(4.5), Inches(4.0), Inches(1.1), PRIMARY)
    add_text(s, Inches(1.0), Inches(4.5), Inches(4.0), Inches(1.1),
             "PostgreSQL 16 (Primary)\nstreaming replication",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round_rect(s, Inches(5.3), Inches(4.5), Inches(3.0), Inches(1.1), GRAY)
    add_text(s, Inches(5.3), Inches(4.5), Inches(3.0), Inches(1.1),
             "PostgreSQL (Replica)\nread-only · tайлан",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round_rect(s, Inches(8.6), Inches(4.5), Inches(3.7), Inches(1.1), GOLD)
    add_text(s, Inches(8.6), Inches(4.5), Inches(3.7), Inches(1.1),
             "Redis (session, cache)\nS3 (images, backup)",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Туслах тайлбар
    add_text(s, Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.0),
             ("CI/CD: GitHub Actions → Docker Hub → SSH deploy · Prometheus + Grafana мониторинг\n"
              "Backup: PostgreSQL pg_dump өдөр бүр, S3-д хадгална · Retention 30 хоног"),
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_deploy)

# -------- Slide 21: Туршилтын төлөвлөгөө --------
def slide_testing():
    s = add_slide(); header(s, "Туршилтын стратеги")
    types = [
        ("Функциональ", "Jest (unit) + Supertest (integration) — 140+ кейс"),
        ("API",         "Postman collection + newman CI-д — 50+ endpoint"),
        ("Ачаалал",     "k6 — 200 VU ачаалалт, p95 < 300мс зорилт"),
        ("Гар (manual)","Test case 60+ (зочин/ажилтан/менежер)"),
        ("Хэрэглэхүй",  "SUS survey, 15 хэрэглэгчтэй туршилт"),
        ("Аюулгүй",     "OWASP ZAP, JWT fuzz, RBAC bypass туршилт"),
    ]
    cw = Inches(4.05); ch = Inches(1.5); gap_x = Inches(0.1); gap_y = Inches(0.15)
    for i, (t, d) in enumerate(types):
        col = i % 3; row = i // 3
        x = Inches(0.5) + col * (cw + gap_x)
        y = Inches(1.3) + row * (ch + gap_y)
        card(s, x, y, cw, ch, t, d, title_color=PRIMARY, bg=LIGHT)
    # Цонх - purity зорилт
    add_round_rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
                   PRIMARY, line_color=PRIMARY)
    add_text(s, Inches(0.7), Inches(4.7), Inches(12.0), Inches(0.45),
             "Чанарын зорилт (Quality Gates)", size=15, bold=True, color=GOLD)
    bullet_list(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.6), [
        ("Code coverage ≥ 70%", "critical route-уудад ≥ 85%"),
        ("API p95 < 300мс",     "/orders, /menu, /tables endpoint-ууд"),
        ("SUS > 75",            "Good — industry benchmark"),
        ("Zero critical CVE",   "npm audit, Trivy SCA"),
    ], size=13, color=WHITE, lead_color=GOLD)

SLIDES.append(slide_testing)

# -------- Slide 22: Ачааллын туршилтын үр дүн --------
def slide_perf():
    s = add_slide(); header(s, "Ачааллын туршилтын үр дүн (k6)")
    # Big numbers row
    stat_card(s, Inches(0.5),  Inches(1.4), Inches(3.0), Inches(2.0),
              "240мс", "API p95 хариу (зорилт < 300)", val_color=GREEN)
    stat_card(s, Inches(3.7),  Inches(1.4), Inches(3.0), Inches(2.0),
              "72мс",  "WebSocket дундаж (зорилт < 100)", val_color=GREEN)
    stat_card(s, Inches(6.9),  Inches(1.4), Inches(3.0), Inches(2.0),
              "0.08%", "Алдааны хувь (зорилт < 1%)", val_color=GREEN)
    stat_card(s, Inches(10.1), Inches(1.4), Inches(3.0), Inches(2.0),
              "220 VU","Тогтвортой ачаалал", val_color=PRIMARY)
    # Хүснэгт
    rows = [
        ("50 VU",  "120мс",  "22мс", "0.01%",  "Ок"),
        ("100 VU", "180мс",  "45мс", "0.02%",  "Ок"),
        ("150 VU", "210мс",  "60мс", "0.04%",  "Ок"),
        ("200 VU", "240мс",  "72мс", "0.08%",  "Ок"),
        ("250 VU", "360мс",  "98мс", "0.4%",   "Болзошгүй"),
        ("300 VU", "620мс",  "140мс","1.8%",   "Хязгаар"),
    ]
    make_table(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(3.4),
               ["VU", "API p95", "WS dundaj", "Алдаа", "Төлөв"], rows,
               col_widths=[2, 3, 3, 2, 3], header_size=13, body_size=12)

SLIDES.append(slide_perf)

# -------- Slide 22b: Ачааллын шугаман график --------
def slide_perf_chart():
    s = add_slide(); header(s, "Ачааллын тест — хариуны хугацааны график")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "k6 ачаалалын тест · virtual user (VU) тоо vs хариуны хугацаа (мс)",
             size=13, italic=True, color=ACCENT)
    cats = ["50", "100", "150", "200", "250", "300"]
    add_chart(
        s, XL_CHART_TYPE.LINE,
        Inches(0.5), Inches(1.7), Inches(8.5), Inches(5.2),
        categories=cats,
        series=[
            ("API p95 (мс)",   (120, 180, 210, 240, 360, 620)),
            ("WebSocket (мс)", (22,  45,  60,  72,  98,  140)),
            ("Зорилт p95 (мс)",(300, 300, 300, 300, 300, 300)),
        ],
        title="Хариуны хугацаа vs ачаалал",
        palette=[PRIMARY, GREEN, RED],
    )
    # Дагалдах тэмдэглэл
    add_round_rect(s, Inches(9.2), Inches(1.7), Inches(3.8), Inches(5.2),
                   LIGHT, line_color=ACCENT)
    add_text(s, Inches(9.4), Inches(1.85), Inches(3.5), Inches(0.4),
             "Үр дүн", size=15, bold=True, color=PRIMARY)
    bullet_list(s, Inches(9.4), Inches(2.3), Inches(3.5), Inches(4.6), [
        ("Тогтвортой", "200 VU хүртэл зорилт хангана."),
        ("Break-point", "≈ 250 VU дээр график огцом өсөх."),
        ("Хязгаар", "300 VU алдаа 1.8% → scale хэрэгтэй."),
        ("WebSocket", "140мс-ээс хэтрэхгүй, push тогтвортой."),
    ], size=12)

SLIDES.append(slide_perf_chart)

# -------- Slide 23: Хэрэглэхүйн үнэлгээ --------
def slide_sus():
    s = add_slide(); header(s, "Хэрэглэхүйн үнэлгээ (SUS)")
    # Том үр дүн
    add_round_rect(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.5),
                   PRIMARY, line_color=PRIMARY)
    add_text(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
             "System Usability Scale", size=14, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.5), Inches(1.85), Inches(5.5), Inches(2.4),
             "82.4", size=96, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(4.2), Inches(5.5), Inches(0.5),
             "Excellent (>80) — industry-ийн дээд гуравны нэгд",
             size=13, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Туршилтын дэлгэрэнгүй
    rows = [
        ("Нийт оролцогч",   "15 хүн"),
        ("Хэрэглэгч",        "10 (5 зочин, 5 хэрэглэгч)"),
        ("Ажилтан",          "3 кассир, 2 менежер"),
        ("Даалгавар",        "6 сценари · дундаж 4:32 минут"),
        ("Амжилттай",        "93.3% (84/90 даалгавар)"),
        ("Алдаа",            "Голчлон QR скан, 2 дахь эргэлтэд 100%"),
        ("Гол санал",        "Монгол мэдэгдэл, хурдтай, ойлгомжтой"),
    ]
    make_table(s, Inches(6.3), Inches(1.3), Inches(6.5), Inches(5.5),
               ["Үзүүлэлт", "Үр дүн"], rows,
               col_widths=[3, 6], header_size=13, body_size=12)

SLIDES.append(slide_sus)

# -------- Slide 24: Бизнесийн үр дүн --------
def slide_business():
    s = add_slide(); header(s, "Бодит нөхцөлд турших үр дүн")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "«Мандуул» ресторан · 3 долоо хоногийн пилот (N=312 хэрэглэгч, 1,480 захиалга)",
             size=13, italic=True, color=ACCENT)
    # Том харьцуулсан үзүүлэлт
    kpi = [
        ("Захиалга өгөх хугацаа", "−62%", "7 мин → 2:40 мин", GREEN),
        ("Захиалгын алдаа",       "−74%", "14% → 3.6%",      GREEN),
        ("Цэсийн шинэчлэлт",      "−95%", "өдөр → шууд",     GREEN),
        ("Хэрэглэгчийн ханамж",   "+41%", "3.4 → 4.8 / 5",   GREEN),
        ("Ажилтны ачаалал",       "−38%", "90 → 56 алхам/ш",GREEN),
        ("Тайлангийн хугацаа",    "−90%", "2 цаг → 12 мин",  GREEN),
    ]
    cw = Inches(4.05); ch = Inches(1.7); gap_x = Inches(0.1); gap_y = Inches(0.15)
    for i, (t, v, sub, c) in enumerate(kpi):
        col = i % 3; row = i // 3
        x = Inches(0.5) + col * (cw + gap_x)
        y = Inches(1.7) + row * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, WHITE, line_color=c)
        add_text(s, x, y + Inches(0.1), cw, Inches(0.4),
                 t, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.5), cw, Inches(0.7),
                 v, size=30, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(1.2), cw, Inches(0.4),
                 sub, size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    # Нэмэгдэл
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2),
             "Өдрийн борлуулалт: +18% · Түүхий эдийн хаягдал: −22% · Хөрөнгө оруулалтын эргэлт (ROI): ~4.5 сар",
             size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

SLIDES.append(slide_business)

# -------- Slide 24b: Before/After харьцуулсан бар график --------
def slide_before_after():
    s = add_slide(); header(s, "Уламжлалт арга vs систем — KPI харьцуулалт")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "«Мандуул» ресторан · 3 долоо хоногийн пилотын бодит хэмжилт",
             size=13, italic=True, color=ACCENT)
    cats = ["Захиалгын\nхугацаа (мин)", "Захиалгын\nалдаа (%)",
            "Цэс\nшинэчлэл (мин)", "Хэрэглэгчийн\nханамж (/5)",
            "Тайлангийн\nхугацаа (мин)"]
    add_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.7), Inches(8.5), Inches(5.2),
        categories=cats,
        series=[
            ("Өмнө (гар аргаар)", (7.0, 14.0, 1440, 3.4, 120)),
            ("Дараа (систем)",    (2.7, 3.6, 1,    4.8, 12)),
        ],
        title="Before vs After — 5 үндсэн үзүүлэлт",
        palette=[RED, GREEN],
    )
    # Summary
    add_round_rect(s, Inches(9.2), Inches(1.7), Inches(3.8), Inches(5.2),
                   LIGHT, line_color=ACCENT)
    add_text(s, Inches(9.4), Inches(1.85), Inches(3.5), Inches(0.4),
             "Нийт сайжруулалт", size=15, bold=True, color=PRIMARY)
    bullet_list(s, Inches(9.4), Inches(2.3), Inches(3.5), Inches(4.6), [
        ("−62%", "Захиалгын хугацаа."),
        ("−74%", "Захиалгын алдаа."),
        ("−95%", "Цэсийн шинэчлэлт."),
        ("+41%", "Хэрэглэгчийн ханамж."),
        ("−90%", "Тайлангийн хугацаа."),
        ("+18%", "Өдрийн борлуулалт."),
    ], size=12, lead_color=GREEN)

SLIDES.append(slide_before_after)

# -------- Slide 24c: SUS хэрэглэхүйн бар график --------
def slide_sus_chart():
    s = add_slide(); header(s, "SUS үнэлгээний дэлгэрэнгүй задаргаа")
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "10 асуултын дундаж үнэлгээ (1–5 шкала · 15 оролцогч)",
             size=13, italic=True, color=ACCENT)
    # SUS 10 questions (Mongolian summary)
    cats = ["Q1\nДахин\nхэрэглэнэ",
            "Q2\nХүндрэлтэй",
            "Q3\nХялбар",
            "Q4\nТусламж\nхэрэгтэй",
            "Q5\nФункц\nнэгдсэн",
            "Q6\nЗөрчилтэй",
            "Q7\nХурдан\nсурна",
            "Q8\nТөвөгтэй",
            "Q9\nИтгэлтэй",
            "Q10\nСурах\nихтэй"]
    add_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2),
        categories=cats,
        series=[("Дундаж үнэлгээ (эерэг чиг)",
                 (4.6, 1.5, 4.7, 1.8, 4.4, 1.6, 4.5, 1.7, 4.3, 1.9))],
        title="SUS нийт оноо = 82.4 / 100 (Excellent)",
        legend=False,
        palette=[PRIMARY],
    )

SLIDES.append(slide_sus_chart)

# -------- Slide 25: Дүгнэлт --------
def slide_conclusion():
    s = add_slide(); header(s, "Дүгнэлт")
    items = [
        ("Бүрэн шийдэл", "Зоогийн газрын ширээ захиалга, цэс удирдлага, төлбөр, тайлангийн бүх мөчлөгийг нэг системд нэгтгэв."),
        ("Монголын онцлог", "QPay, монгол хэл, орон нутгийн VPS — 95% хямд, хэрэглэгчид танил."),
        ("Хэмжигдэхүйц үр дүн", "Захиалгын хугацаа −62%, алдаа −74%, SUS 82.4 (excellent)."),
        ("Архитектур", "Layered + monorepo + Docker — 1 хөгжүүлэгчээс 10 хөгжүүлэгч багт ажлын бэлэн бүтэц."),
        ("Технологийн баталгаа", "12 FR, 25+ NFR, 30 API, 9 Socket эвэнт, 9 үндсэн entity — бүгд туршигдсан."),
        ("Бизнес утга", "ROI ≈ 4.5 сар, борлуулалт +18% — Монголын SME-д хүртээмжтэй цифр шилжилт."),
    ]
    bullet_list(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7), items,
                size=15, lead_color=GREEN)

SLIDES.append(slide_conclusion)

# -------- Slide 26: Цаашдын ажил --------
def slide_future():
    s = add_slide(); header(s, "Цаашдын судалгааны чиглэл")
    items = [
        ("Олон салбарт", "Нэг толгой компани, олон ресторан — мульти-тенант панель, тусгай эрхийн бүтэц."),
        ("AR цэс",        "Хоолны 3D загварчлал, AR-ээр харах боломж (ARKit / ARCore)."),
        ("ML санал",      "Захиалгын түүхэнд суурилсан персонал санал (collaborative filtering)."),
        ("Voice захиалга","Mongolian ASR интеграц — ярианы захиалга, ахмадуудад зориулсан."),
        ("Delivery",      "Хүргэлт, жолоочийн апп, газрын зурагт дагалдан хянах."),
        ("Нэгдсэн ERP",   "Нягтлан бодох, цалин, хангамжийн гүн интеграц."),
        ("Глобал",        "Stripe, Apple/Google Pay, Англи+Хятад хэл, дэлхийн зах зээлд."),
    ]
    bullet_list(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7), items,
                size=14, lead_color=GOLD)

SLIDES.append(slide_future)

# -------- Slide 27: Талархал / Асуулт --------
def slide_thanks():
    s = add_slide(); set_bg(s, PRIMARY)
    add_rect(s, 0, 0, SW, Inches(0.15), GOLD)
    add_rect(s, 0, SH - Inches(0.15), SW, Inches(0.15), GOLD)
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
             "АНХААРАЛ\nХАНДУУЛСАНД\nБАЯРЛАЛАА",
             size=56, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), GOLD)
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             "Асуулт хэлэлцүүлэг",
             size=22, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
             "Эрдэнэтөгсийн Очбадрах",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
             "sergelenbayar0@gmail.com",
             size=13, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "Өмнөговь аймаг дахь Технологийн Дээд Сургууль · 2026",
             size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

SLIDES.append(slide_thanks)

# ======================= Үүсгэх =======================
for i, fn in enumerate(SLIDES):
    fn()

# Footer бүх слайдад (эхний болон сүүлийн cover/thanks-ийг алгасна)
total = len(prs.slides)
for idx, slide in enumerate(prs.slides, start=1):
    # Cover (1) ба Thanks (total) өнгөтэй тул footer нэмэхгүй
    if idx == 1 or idx == total:
        continue
    footer(slide, idx, total)

out = r"c:\Users\ochko\Downloads\Restaurant-Table-Booking-1\Restaurant-Table-Booking-1\artifacts\presentation\Restaurant-Booking-Thesis.pptx"
prs.save(out)
print(f"OK: {total} slides -> {out}")
