# -*- coding: utf-8 -*-
"""
QR-кодон Ресторан Захиалгын Систем
Бакалаврын дипломын ажлын хамгаалалтын илтгэл — Системийн АЖИЛЛАХ ЯВЦ
зориулсан 14-слайдын хувилбар (16:9). 10-15 минутын ярианд тохирсон.

Гаралт: workflow.pptx
"""

from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────── Өнгөний схем ──────────────────────────
PRIMARY = RGBColor(0x00, 0x46, 0x7F)
ACCENT  = RGBColor(0x46, 0x82, 0xB4)
GOLD    = RGBColor(0xD4, 0xA0, 0x17)
DARK    = RGBColor(0x1E, 0x1E, 0x1E)
GRAY    = RGBColor(0x60, 0x60, 0x60)
LIGHT   = RGBColor(0xF2, 0xF6, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
RED     = RGBColor(0xC6, 0x28, 0x28)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
PURPLE  = RGBColor(0x6A, 0x1B, 0x9A)
TEAL    = RGBColor(0x00, 0x8B, 0x8B)

FONT = "Calibri"
TOTAL_SLIDES = 16

# Дипломын зургуудын зам — slide 5, 6, 6b-д жинхэнэ screenshot ашиглах
FIG_DIR = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "api-server", "thesis", "Figures",
)

def fig(name: str) -> str:
    return os.path.join(FIG_DIR, name)

def add_picture_fit(slide, path, x, y, w, h, *, fallback_text="[Зураг олдсонгүй]"):
    """Зургийг өгсөн хайрцагт fit хийж нэмнэ; олдохгүй бол placeholder зурдаг."""
    if os.path.exists(path):
        # Хайрцагтай ижил пропорцоор тааруулж зурдаг учир shrink-to-fit
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    add_round_rect(slide, x, y, w, h, LIGHT, line_color=GRAY)
    add_text(slide, x, y, w, h, fallback_text, size=11, italic=True,
             color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return None

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ─────────────────────── Туслах функцууд ───────────────────────
def add_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s

def add_round_rect(slide, x, y, w, h, color, line_color=None, line_pt=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.05), GOLD)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11.5), Inches(0.7),
             title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.8), Inches(0.18), Inches(2.3), Inches(0.35),
             "Өмнөговь ТДС", size=11, color=WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.8), Inches(0.5), Inches(2.3), Inches(0.3),
             "Бакалаврын дипломын ажил", size=9, italic=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
                 subtitle, size=13, italic=True, color=ACCENT)

def footer(slide, page_num, total=TOTAL_SLIDES):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), PRIMARY)
    add_text(slide, Inches(0.4), SH - Inches(0.32), Inches(9), Inches(0.32),
             "QR-кодон ресторан захиалгын систем — Системийн ажиллах явц",
             size=9, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2.2), SH - Inches(0.32), Inches(2.0), Inches(0.32),
             f"{page_num} / {total}", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def bullet_list(slide, x, y, w, h, items, *, size=14, color=DARK,
                lead_color=PRIMARY, line_spacing=1.2, bullet="▸"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r0 = p.add_run()
        r0.text = f"{bullet}  "
        r0.font.name = FONT; r0.font.size = Pt(size)
        r0.font.bold = True; r0.font.color.rgb = lead_color
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.name = FONT; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = lead_color
            r2 = p.add_run()
            r2.text = " — " + rest
            r2.font.name = FONT; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r1 = p.add_run()
            r1.text = str(item)
            r1.font.name = FONT; r1.font.size = Pt(size)
            r1.font.color.rgb = color
    return tb

def card(slide, x, y, w, h, title, body, *, title_color=PRIMARY, bg=LIGHT,
         title_size=14, body_size=11, icon=None):
    add_round_rect(slide, x, y, w, h, bg, line_color=ACCENT)
    top = y + Inches(0.1)
    if icon:
        add_text(slide, x + Inches(0.2), top, Inches(0.5), Inches(0.5),
                 icon, size=22, bold=True, color=title_color,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.7), top, w - Inches(0.85), Inches(0.5),
                 title, size=title_size, bold=True, color=title_color,
                 anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_text(slide, x + Inches(0.15), top, w - Inches(0.3), Inches(0.45),
                 title, size=title_size, bold=True, color=title_color)
    add_text(slide, x + Inches(0.15), y + Inches(0.65), w - Inches(0.3),
             h - Inches(0.75), body, size=body_size, color=DARK,
             line_spacing=1.2)

def make_table(slide, x, y, w, h, headers, rows, *,
               header_color=PRIMARY, zebra=True, header_size=11, body_size=10,
               col_widths=None):
    cols = len(headers); rws = len(rows) + 1
    tbl_shape = slide.shapes.add_table(rws, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for c, h_ in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(54000); tf.margin_right = Emu(54000)
        tf.margin_top = Emu(36000); tf.margin_bottom = Emu(36000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for r in list(p.runs): r._r.getparent().remove(r._r)
        run = p.add_run(); run.text = h_
        run.font.name = FONT; run.font.size = Pt(header_size)
        run.font.bold = True; run.font.color.rgb = WHITE
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            if zebra and ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = Emu(54000); tf.margin_right = Emu(54000)
            tf.margin_top = Emu(27000); tf.margin_bottom = Emu(27000)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            for r in list(p.runs): r._r.getparent().remove(r._r)
            run = p.add_run(); run.text = str(val)
            run.font.name = FONT; run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
    return tbl

def arrow(slide, x1, y1, x2, y2, color=PRIMARY, pt=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(pt)
    return line

def phone_mockup(slide, x, y, w, h, *, title="Screen", lines=None, bg=DARK,
                 inner_bg=WHITE):
    """Утасны хэлбэртэй mockup зурдаг (frame + screen + content)."""
    # Outer frame
    add_round_rect(slide, x, y, w, h, bg, line_color=DARK, line_pt=2)
    # Notch
    notch_w = Inches(0.6); notch_h = Inches(0.12)
    add_round_rect(slide, x + (w - notch_w) / 2, y + Inches(0.08),
                   notch_w, notch_h, DARK)
    # Screen area
    sx = x + Inches(0.12); sy = y + Inches(0.32)
    sw = w - Inches(0.24); sh = h - Inches(0.5)
    add_rect(slide, sx, sy, sw, sh, inner_bg)
    # Title bar
    add_rect(slide, sx, sy, sw, Inches(0.4), PRIMARY)
    add_text(slide, sx, sy + Inches(0.05), sw, Inches(0.3),
             title, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body lines
    if lines:
        for i, line in enumerate(lines):
            ly = sy + Inches(0.5) + Inches(0.35 * i)
            if ly + Inches(0.3) > sy + sh:
                break
            if isinstance(line, tuple):
                txt, color = line
            else:
                txt, color = line, DARK
            add_text(slide, sx + Inches(0.1), ly, sw - Inches(0.2), Inches(0.3),
                     txt, size=9, color=color, anchor=MSO_ANCHOR.MIDDLE)

def desktop_mockup(slide, x, y, w, h, *, title="Dashboard", sections=None):
    """Ажилтны самбарын mockup (top bar + sidebar + content)."""
    add_round_rect(slide, x, y, w, h, DARK, line_color=DARK, line_pt=2)
    # Top bar
    add_rect(slide, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2),
             Inches(0.45), PRIMARY)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.5),
             Inches(0.35), title, size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Content background
    add_rect(slide, x + Inches(0.1), y + Inches(0.6), w - Inches(0.2),
             h - Inches(0.7), WHITE)
    if sections:
        sx = x + Inches(0.2); sy = y + Inches(0.75)
        sw = w - Inches(0.4)
        for i, (h_, body, color) in enumerate(sections):
            ty = sy + Inches(0.5 * i)
            if ty + Inches(0.5) > y + h - Inches(0.1):
                break
            add_round_rect(slide, sx, ty, sw, Inches(0.42), LIGHT,
                           line_color=color)
            add_text(slide, sx + Inches(0.1), ty, Inches(2.2), Inches(0.42),
                     h_, size=10, bold=True, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, sx + Inches(2.4), ty, sw - Inches(2.5), Inches(0.42),
                     body, size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════
#  SLIDE 1 — Гарчиг
# ════════════════════════════════════════════════════════════════
def slide_1_title():
    s = add_slide()
    set_bg(s, WHITE)
    add_rect(s, 0, 0, Inches(0.5), SH, PRIMARY)
    add_rect(s, Inches(0.5), 0, Inches(0.15), SH, GOLD)

    add_text(s, Inches(1.0), Inches(0.8), Inches(11), Inches(0.5),
             "ӨМНӨГОВЬ АЙМАГ ДАХЬ ТЕХНОЛОГИЙН ДЭЭД СУРГУУЛЬ",
             size=18, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(1.35), Inches(11), Inches(0.35),
             "Програм хангамжийн тэнхим",
             size=14, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.04), GOLD)

    add_text(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.0),
             "QR-кодоор ширээнээс шууд хоол захиалах",
             size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.9),
             "бодит цаг хугацааны вэб систем",
             size=28, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.4),
             "— Системийн ажиллах явц —",
             size=16, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

    box_x = Inches(3.0); box_y = Inches(5.0)
    box_w = Inches(7.3); box_h = Inches(1.7)
    add_round_rect(s, box_x, box_y, box_w, box_h, LIGHT, line_color=ACCENT)
    add_text(s, box_x + Inches(0.3), box_y + Inches(0.15), box_w - Inches(0.6),
             Inches(0.4), "Зохиогч:  Эрдэнэтөгсийн Очбадрах",
             size=14, bold=True, color=DARK)
    add_text(s, box_x + Inches(0.3), box_y + Inches(0.55), box_w - Inches(0.6),
             Inches(0.4), "Удирдагч:  П. Нарангэрэл",
             size=12, color=DARK)
    add_text(s, box_x + Inches(0.3), box_y + Inches(0.9), box_w - Inches(0.6),
             Inches(0.4), "Зөвлөгч:  доктор Б. Намжилдорж",
             size=12, color=DARK)
    add_text(s, box_x + Inches(0.3), box_y + Inches(1.25), box_w - Inches(0.6),
             Inches(0.4),
             "Огноо:  2026 он, 5-р сар   ·   Даланзадгад, Өмнөговь",
             size=11, italic=True, color=GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 2 — Үндэслэл
# ════════════════════════════════════════════════════════════════
def slide_2_rationale():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Үндэслэл", "Уламжлалт үйлчилгээний урсгал ба түүний дутагдал")

    # Left: traditional flow problems
    left_x = Inches(0.6); col_w = Inches(6.0)
    add_text(s, left_x, Inches(1.55), col_w, Inches(0.45),
             "Уламжлалт урсгалд тулгардаг саад:",
             size=15, bold=True, color=PRIMARY)
    problems = [
        ("Хүлээлт", "Үйлчлэгч олон ширээ хооронд гүйж хүрэхгүй, зочин 10-15 мин хүлээдэг"),
        ("Алдаа", "Цаасан дээр гар бичих, аман дамжуулах нь буруу захиалга үүсгэдэг"),
        ("Харьцалт", "Менежер хоол, нөөц, орлогыг бодит цагт харах боломжгүй"),
        ("Олон зочин", "Нэг ширээний хүн бүрд тус тус үйлчлэх — олон удаа эргэнэ"),
        ("Хэлний саад", "Гадаад зочдод үйлчлэгч цэс тайлбарлахад бэрхшээлтэй"),
    ]
    bullet_list(s, left_x, Inches(2.05), col_w, Inches(4.8), problems,
                size=13, line_spacing=1.3)

    # Right: market data + opportunity
    right_x = Inches(7.0); col_w2 = Inches(5.8)
    add_text(s, right_x, Inches(1.55), col_w2, Inches(0.45),
             "Зах зээлийн боломж (Монгол улс):",
             size=15, bold=True, color=ORANGE)

    stats = [
        ("78.4%", "Смартфоны хүртээмж"),
        ("92%", "Шинэ технологи нэвтрүүлэхэд саад туларсан рестораны хувь"),
        ("5,200+", "Улаанбаатар хотын хүнсний үйлчилгээний газар"),
        ("185%", "QR ашиглалт 2020-2023 онд өссөн (Ази)"),
    ]
    sy = Inches(2.1)
    for i, (val, label) in enumerate(stats):
        y = sy + Inches(1.1 * i)
        add_round_rect(s, right_x, y, Inches(1.6), Inches(0.95), PRIMARY)
        add_text(s, right_x, y + Inches(0.05), Inches(1.6), Inches(0.55),
                 val, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, right_x, y + Inches(0.55), Inches(1.6), Inches(0.4),
                 "stat", size=8, color=GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, right_x + Inches(1.8), y + Inches(0.15),
                 col_w2 - Inches(1.9), Inches(0.65),
                 label, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.15)

    footer(s, 2)


# ════════════════════════════════════════════════════════════════
#  SLIDE 3 — Ижил төстэй системийн харьцуулалт
# ════════════════════════════════════════════════════════════════
def slide_3_comparison():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Ижил төстэй системийн харьцуулсан судалгаа",
           "Гадаад худалдааны системүүд vs Энэхүү шийдэл")

    headers_row = ["Шалгуур", "Toast POS\n(АНУ)", "Lightspeed\n(Канад)",
                   "Square\n(АНУ)", "Энэхүү систем"]
    rows = [
        ["QR захиалга",           "Нэмэлт",   "Нэмэлт",   "Тийм",     "Гол функц"],
        ["Realtime WebSocket",    "Хэсэгчлэн","Тийм",     "Хэсэгчлэн","Тийм"],
        ["Multi-guest нэг ширээнд","Үгүй",    "Үгүй",     "Үгүй",     "Тийм"],
        ["Монгол хэл",            "Үгүй",     "Үгүй",     "Үгүй",     "Тийм"],
        ["QPay / SocialPay",      "Үгүй",     "Үгүй",     "Үгүй",     "Бэлэн"],
        ["Сарын зардал (USD)",    "69-165",   "69-135",   "60+",      "~10-25"],
        ["AR 3D загвар",          "Үгүй",     "Үгүй",     "Үгүй",     "Тийм"],
        ["Нээлттэй эх код",       "Үгүй",     "Үгүй",     "Үгүй",     "Тийм"],
    ]
    make_table(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(4.8),
               headers_row, rows,
               header_size=11, body_size=11,
               col_widths=[2.5, 1.6, 1.6, 1.6, 2.0])

    # Note
    add_round_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.55),
                   LIGHT, line_color=ORANGE)
    add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.55),
             "Энэхүү системийн давуу тал: open-source, монгол хэлтэй, локал "
             "төлбөртэй интегралд бэлэн, multi-guest + AR-тэй цорын ганц шийдэл.",
             size=11, italic=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 3)


# ════════════════════════════════════════════════════════════════
#  SLIDE 4 — Зорилго + Архитектур (товч)
# ════════════════════════════════════════════════════════════════
def slide_4_goal_arch():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Зорилго ба архитектур")

    # Goal box
    add_round_rect(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.1), PRIMARY)
    add_text(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(0.4),
             "ГОЛ ЗОРИЛГО", size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(0.6),
             "QR кодоор зочин өөрөө захиалга өгч, ажилтан real-time хүлээн авч "
             "боловсруулдаг, менежер нөөц, орлогоо хянадаг цогц вэб систем.",
             size=13, color=WHITE, line_spacing=1.25)

    # 4-layer architecture
    add_text(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(0.4),
             "4 давхаргат архитектур:",
             size=14, bold=True, color=PRIMARY)
    layers = [
        ("КЛИЕНТ", "React 19 · Vite 7 · TanStack Query · Tailwind CSS · shadcn/ui",
         RGBColor(0xDC, 0xE9, 0xF5), PRIMARY),
        ("API", "Express 5 · TypeScript · JWT · 43 REST endpoint · zod validation",
         RGBColor(0xFD, 0xE7, 0xC9), ORANGE),
        ("REALTIME", "Socket.IO 4 · Rooms · 12 event төрөл · авто-reconnect",
         RGBColor(0xD8, 0xF1, 0xD8), GREEN),
        ("ӨГӨГДӨЛ", "PostgreSQL 16 · Drizzle ORM · 14 хүснэгт · ACID транзакц",
         RGBColor(0xE9, 0xD8, 0xF1), PURPLE),
    ]
    ly = Inches(3.4); lh = Inches(0.85)
    for i, (h_, t, bg, c) in enumerate(layers):
        y = ly + i * (lh + Inches(0.1))
        add_round_rect(s, Inches(0.6), y, Inches(12.1), lh, bg, line_color=c)
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(2.4), Inches(0.6),
                 h_, size=15, bold=True, color=c,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.3), y + Inches(0.1), Inches(9.3), Inches(0.6),
                 t, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 4)


# ════════════════════════════════════════════════════════════════
#  SLIDE 5 — Screenshot #1: Public + Guest
# ════════════════════════════════════════════════════════════════
def slide_5_screens_public():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Системийн дэлгэцүүд (1) — Зочны хэсэг",
           "Гар утсаар QR-ийг уншуулсны дараах зочны 3 үндсэн дэлгэц")

    # 3 mobile screenshots side by side (16:9 slide)
    phone_w = Inches(2.6); phone_h = Inches(5.0)
    px = [Inches(1.2), Inches(5.35), Inches(9.5)]; py = Inches(1.55)
    captions = [
        ("guest-menu.png",         "Зочны цэс  (/menu?t=...)"),
        ("guest-cart.png",         "Сагс  (/cart)"),
        ("guest-order-status.png", "Захиалгын төлөв  (/order-status)"),
    ]
    for i, (fname, cap) in enumerate(captions):
        add_picture_fit(s, fig(fname), px[i], py, phone_w, phone_h,
                        fallback_text=f"[{fname}]")
        add_text(s, px[i], py + phone_h + Inches(0.1), phone_w, Inches(0.35),
                 cap, size=11, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER)

    footer(s, 5)


# ════════════════════════════════════════════════════════════════
#  SLIDE 6 — Screenshot #2: Staff dashboard
# ════════════════════════════════════════════════════════════════
def slide_6_screens_staff():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Системийн дэлгэцүүд (2) — Ажилтны самбар",
           "Захиалгын урсгал ба ширээний удирдлагын дэлгэц")

    # 2 desktop screenshots side by side
    img_w = Inches(6.1); img_h = Inches(4.4)
    add_picture_fit(s, fig("staff-dashboard.png"),
                    Inches(0.35), Inches(1.7), img_w, img_h,
                    fallback_text="[staff-dashboard.png]")
    add_text(s, Inches(0.35), Inches(6.15), img_w, Inches(0.3),
             "Захиалгын самбар  (/admin)", size=11, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)

    add_picture_fit(s, fig("staff-tables.png"),
                    Inches(6.85), Inches(1.7), img_w, img_h,
                    fallback_text="[staff-tables.png]")
    add_text(s, Inches(6.85), Inches(6.15), img_w, Inches(0.3),
             "Ширээ + QR код  (/admin/tables)", size=11, bold=True,
             color=PRIMARY, align=PP_ALIGN.CENTER)

    footer(s, 6)


# ════════════════════════════════════════════════════════════════
#  SLIDE 6b — Screenshot #3: Цэс, тайлан, нэвтрэх
# ════════════════════════════════════════════════════════════════
def slide_6b_screens_admin_more():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Системийн дэлгэцүүд (3) — Цэс, тайлан, нэвтрэх",
           "Менежерийн нэмэлт удирдлагын модулиуд")

    img_w = Inches(4.05); img_h = Inches(4.3)
    px = [Inches(0.35), Inches(4.65), Inches(8.95)]
    py = Inches(1.7)
    captions = [
        ("login-page.png",        "Нэвтрэх  (/login)"),
        ("staff-menu-manage.png", "Цэс удирдлага  (/admin/menu)"),
        ("staff-reports.png",     "Тайлан  (/admin/reports)"),
    ]
    for i, (fname, cap) in enumerate(captions):
        add_picture_fit(s, fig(fname), px[i], py, img_w, img_h,
                        fallback_text=f"[{fname}]")
        add_text(s, px[i], py + img_h + Inches(0.05), img_w, Inches(0.3),
                 cap, size=11, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER)

    footer(s, 7)


# ════════════════════════════════════════════════════════════════
#  SLIDE 7 — Захиалгын ерөнхий ажиллах урсгал
# ════════════════════════════════════════════════════════════════
def slide_7_workflow():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Зочны захиалгын бүрэн ажиллах урсгал",
           "QR scan-аас төлбөр баталгаажтал явах 7 алхам")

    steps = [
        ("📱", "1. QR scan",     "Зочин ширээний QR-ыг утсаар уншуулна"),
        ("🔗", "2. Session",     "Сервер ширээний qr_token шалгаж\nsession JWT олгоно"),
        ("🍴", "3. Цэс",          "Цэснээс хоол сонгож сагсанд оруулна\n(AR view-аар үзэх боломжтой)"),
        ("📤", "4. Илгээх",       "POST /api/cart → DB-д захиалга үүснэ\ninventory автомат хасагдана"),
        ("⚡",  "5. Realtime",   "Socket.IO order:new event\nbroadcast → бүх ажилтан мэдэгдэнэ"),
        ("👨‍🍳","6. Бэлдэлт",    "Chef Kanban-аас 'Бэлдэх' → 'Бэлэн'\nbadge real-time зочны утсанд"),
        ("💳", "7. Төлбөр",       "Cashier QPay/cash/card сонгоно\nsession хаагдаж ширээ available"),
    ]
    n = len(steps)
    sx = Inches(0.5); sy = Inches(1.7)
    total_w = Inches(12.3)
    gap = Inches(0.15)
    box_w = (total_w - gap * (n - 1)) / n
    box_h = Inches(3.5)
    for i, (icon, title, body) in enumerate(steps):
        x = sx + i * (box_w + gap)
        add_round_rect(s, x, sy, box_w, box_h, LIGHT,
                       line_color=PRIMARY, line_pt=1.5)
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + (box_w - Inches(0.55)) / 2,
                                    sy + Inches(0.15),
                                    Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        add_text(s, x, sy + Inches(0.15), box_w, Inches(0.55),
                 icon, size=22, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, sy + Inches(0.8), box_w, Inches(0.4),
                 title, size=11, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.05), sy + Inches(1.25),
                 box_w - Inches(0.1), Inches(2.1),
                 body, size=9, color=GRAY, align=PP_ALIGN.CENTER,
                 line_spacing=1.2)
        if i < n - 1:
            ax1 = x + box_w
            ax2 = x + box_w + gap
            ay = sy + box_h / 2
            arrow(s, ax1, ay, ax2, ay, PRIMARY, 2)

    # Highlight box
    note_y = sy + box_h + Inches(0.25)
    add_round_rect(s, Inches(0.5), note_y, Inches(12.3), Inches(0.7),
                   GOLD)
    add_text(s, Inches(0.7), note_y, Inches(11.9), Inches(0.7),
             "🔑  Энэ урсгалын 5-6 алхамд хүний оролцоо БАЙХГҮЙ — "
             "захиалга, төлөв шинэчлэлт бүгд автомат + real-time.",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 8)


# ════════════════════════════════════════════════════════════════
#  SLIDE 8 — Multi-guest QR session
# ════════════════════════════════════════════════════════════════
def slide_8_multiguest():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Multi-guest QR session — олон зочин нэг ширээн дээр",
           "Гэр бүлийн, нөхдийн орхиол үед хэрэгцээтэй боломж")

    # Left: 3 phone mockups joining the same session
    px = [Inches(0.6), Inches(2.8), Inches(5.0)]
    py = Inches(1.7); phone_w = Inches(2.0); phone_h = Inches(3.6)
    persons = [("Болд", "🍴 Бууз x2\n🥤 Кофе x1"),
               ("Сараа", "🍜 Цуйван x1\n🧃 Жүүс x1"),
               ("Зул", "🥩 Шарсан мах\n🍰 Бялуу x2")]
    for i, ((name, items), x) in enumerate(zip(persons, px)):
        phone_mockup(s, x, py, phone_w, phone_h,
                     title=f"Ширээ #3 — {name}",
                     lines=[
                         (f"Зочин: {name}", PRIMARY),
                         ("session #12", GRAY),
                         ("", DARK),
                         ("Миний захиалга:", DARK),
                         *[(line, DARK) for line in items.split("\n")],
                         ("", DARK),
                         ("Бусад зочин (2):", ORANGE),
                         ("✓ нэгдсэн", GREEN),
                     ])

    # Arrow to "Single session"
    cx = Inches(7.4); cy = Inches(1.9); cw = Inches(5.4); ch = Inches(3.2)
    add_round_rect(s, cx, cy, cw, ch, PRIMARY)
    add_text(s, cx, cy + Inches(0.15), cw, Inches(0.5),
             "🪑  ШИРЭЭ #3", size=18, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx, cy + Inches(0.7), cw, Inches(0.4),
             "Session #12 (нэгдсэн)",
             size=12, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    inner = [
        ("Болд:",  "Бууз x2 + Кофе x1"),
        ("Сараа:", "Цуйван x1 + Жүүс x1"),
        ("Зул:",   "Шарсан мах + Бялуу x2"),
        ("Бүгд:",  "8 мөр  ·  84,500₮"),
    ]
    iy = cy + Inches(1.15)
    for nm, txt in inner:
        add_text(s, cx + Inches(0.3), iy, Inches(1.5), Inches(0.4),
                 nm, size=12, bold=True, color=GOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(1.5), iy, cw - Inches(1.7), Inches(0.4),
                 txt, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        iy += Inches(0.4)

    # Arrows from phones to central session
    for x in px:
        ax = x + phone_w / 2
        arrow(s, ax + Inches(0.3), Inches(5.4), cx + Inches(0.2),
              cy + ch - Inches(0.3), GOLD, 2)

    # Bottom: how it works (technical)
    by = Inches(5.55)
    add_round_rect(s, Inches(0.6), by, Inches(12.1), Inches(1.55),
                   LIGHT, line_color=ACCENT)
    add_text(s, Inches(0.85), by + Inches(0.1), Inches(11.6), Inches(0.4),
             "🔧 Хэрхэн ажилладаг вэ?",
             size=13, bold=True, color=PRIMARY)
    tech = [
        "1. QR scan-д table_sessions хүснэгтэд нэг session үүсгэнэ "
        "(хэрэв байхгүй бол)",
        "2. Нэмэлт зочин session_participants-д бүртгэгдэж "
        "Socket.IO 'participant:joined' event broadcast хийнэ",
        "3. shared_cart_items хүснэгт нь session-д харьяалагдан хамтын сагсыг "
        "real-time бүх төхөөрөмжид синхрончилно",
        "4. 10 минутын дотор давтан захиалга өгсөн зочны item хуучин "
        "захиалга руу coalescing хийгдэнэ",
    ]
    bullet_list(s, Inches(0.85), by + Inches(0.5), Inches(11.6), Inches(1.05),
                tech, size=10, lead_color=PRIMARY, bullet="•",
                line_spacing=1.15)

    footer(s, 9)


# ════════════════════════════════════════════════════════════════
#  SLIDE 9 — Realtime захиалгын дамжуулалт
# ════════════════════════════════════════════════════════════════
def slide_9_realtime():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Бодит цаг хугацааны захиалгын дамжуулалт",
           "Socket.IO room архитектураар захиалга ширээнээс хүн бүрд")

    # Left: sequence
    add_text(s, Inches(0.6), Inches(1.65), Inches(6.5), Inches(0.4),
             "📡 Захиалгын урсгал (sequence):",
             size=14, bold=True, color=PRIMARY)
    seq = [
        ("Зочин",         "POST /api/cart илгээнэ", PRIMARY),
        ("API сервер",    "DB-д захиалга үүсгэнэ\n(transactional)", ORANGE),
        ("Socket.IO",     "io.to('restaurant_1').emit('order:new')", GREEN),
        ("Chef клиент",   "Kanban-д шинэ карт нэмэгдэнэ", PURPLE),
        ("Chef",          "'Бэлдэх' товч → PATCH /orders/:id/status", ORANGE),
        ("Socket.IO",     "io.to('table_3').emit('order:status')", GREEN),
        ("Зочны клиент",  "Badge real-time шинэчлэгдэнэ", PRIMARY),
    ]
    sy = Inches(2.1)
    for i, (actor, action, color) in enumerate(seq):
        y = sy + i * Inches(0.6)
        # Step circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y,
                                    Inches(0.4), Inches(0.4))
        circle.fill.solid(); circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(s, Inches(0.7), y, Inches(0.4), Inches(0.4),
                 str(i + 1), size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.25), y, Inches(1.7), Inches(0.4),
                 actor, size=11, bold=True, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.0), y, Inches(4.0), Inches(0.6),
                 action, size=10, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas")

    # Right: Room architecture
    rx = Inches(7.5)
    add_text(s, rx, Inches(1.65), Inches(5.3), Inches(0.4),
             "🏠 Socket.IO room архитектур:",
             size=14, bold=True, color=PRIMARY)
    rooms = [
        ("restaurant_1",       "Бүх ажилтан\norder:new, inventory:low, review:new",
         PRIMARY),
        ("session_<qrToken>",  "Тухайн ширээний бүх зочин\nparticipant:joined, cart:updated",
         ORANGE),
        ("table_<tableId>",    "Тухайн ширээний бүх захиалга\norder:status, payment:received",
         GREEN),
    ]
    ry = Inches(2.15)
    for i, (room, audience, color) in enumerate(rooms):
        y = ry + i * Inches(1.3)
        add_round_rect(s, rx, y, Inches(5.3), Inches(1.15),
                       LIGHT, line_color=color)
        add_text(s, rx + Inches(0.15), y + Inches(0.1),
                 Inches(5.0), Inches(0.4),
                 room, size=12, bold=True, color=color, font="Consolas")
        add_text(s, rx + Inches(0.15), y + Inches(0.5),
                 Inches(5.0), Inches(0.65),
                 audience, size=10, color=DARK, line_spacing=1.2)

    # Performance metric
    py_m = Inches(6.3)
    add_round_rect(s, Inches(0.6), py_m, Inches(12.1), Inches(0.7),
                   GOLD)
    add_text(s, Inches(0.85), py_m, Inches(11.7), Inches(0.7),
             "⚡  Хэмжсэн дундаж latency: loopback ~8.5 мс, "
             "Cloudflare tunnel ~142 мс (P95 195 мс).  ",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 10)


# ════════════════════════════════════════════════════════════════
#  SLIDE 10 — Менежерийн боломжууд
# ════════════════════════════════════════════════════════════════
def slide_10_manager():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Менежерийн удирдлагын боломжууд",
           "9 үндсэн модуль — manager role-д бүгд хандах эрхтэй")

    features = [
        ("📊", "Орлогын dashboard",
         "Өдөр / 7 хоног / сар-аар Recharts pie + bar\nТөлбөрийн арга нягтаршил"),
        ("🍴", "Цэс удирдлага",
         "Хоол / ангилал CRUD\nЗураг upload (multer)\n3D model URL"),
        ("🪑", "Ширээ + QR",
         "QR код PDF хэвлэх\nQR rotation\nСтатус (available / occupied)"),
        ("📦", "Нөөц",
         "Inventory CRUD\nАвтомат хасалт\nДоод хэмжээ сэрэмжлүүлэг"),
        ("👥", "Ажилтан + ээлж",
         "Хэрэглэгч CRUD\nRole assign\nShift open/close + кассын нэгтгэл"),
        ("🎯", "Banner / урамшуулал",
         "Зураг + linkUrl нэмэх\nКарусел дараалал\nЭх/идэвхгүй sонголт"),
        ("⭐", "Сэтгэгдэл moderation",
         "1-5 одтой үнэлгээ\nApprove / reject\nReply бичих"),
        ("📅", "Урьдчилсан захиалга",
         "Огноо / цаг / хүний тоо\nManagerт notification\nAvailable шилжүүлэх"),
        ("📋", "Тайлан + орлого",
         "Шилдэг борлуулалт\nНөөц зарцуулалт\nDate range query"),
    ]
    cols = 3; rows = 3
    cw = Inches(4.05); ch = Inches(1.55)
    sx = Inches(0.5); sy = Inches(1.7)
    gx = Inches(0.15); gy = Inches(0.12)
    for i, (icon, title, body) in enumerate(features):
        c = i % cols; r = i // cols
        x = sx + c * (cw + gx)
        y = sy + r * (ch + gy)
        card(s, x, y, cw, ch, title, body, icon=icon,
             title_size=12, body_size=10)

    footer(s, 11)


# ════════════════════════════════════════════════════════════════
#  SLIDE 11 — AR ашиглалтын давуу тал
# ════════════════════════════════════════════════════════════════
def slide_11_ar_benefits():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "AR (Augmented Reality) — давуу тал",
           "Хоолны 3D загварыг утаснаас бодит ширээний дээр харах")

    # Left: benefits
    benefits = [
        ("Захиалгын баталгаа",
         "Зочин хоолыг үзээгүйгээс татгалзах магадлал 35%-иар буурдаг"),
        ("Хоолны хэмжээ",
         "Бодит хэмжээ, тавиурын тоо, гарнитур харах боломж"),
        ("Гадаад зочид",
         "Хоолны нэр ойлгомжгүй ч 3D загвараар сонгож чадна"),
        ("Маркетинг",
         "Шинэ хоол санал болгох — зураг + тайлбараас илүү үр дүнтэй"),
        ("Брендийн ялгарал",
         "Монголд анх удаа AR-тэй ресторан — өрсөлдөөнт давуу тал"),
        ("Хэрэглээ хялбар",
         "Тусдаа app суулгахгүй — браузерт суурилсан"),
    ]
    add_text(s, Inches(0.6), Inches(1.65), Inches(7.5), Inches(0.4),
             "Хэрэглэгчид өгөх давуу тал:",
             size=14, bold=True, color=PRIMARY)
    bullet_list(s, Inches(0.6), Inches(2.15), Inches(7.5), Inches(4.8),
                benefits, size=12, line_spacing=1.3)

    # Right: tech stack
    rx = Inches(8.4)
    add_text(s, rx, Inches(1.65), Inches(4.5), Inches(0.4),
             "Технологийн стек:",
             size=14, bold=True, color=ORANGE)
    tech_items = [
        ("<model-viewer>", "Google web component"),
        (".glb формат", "GLTF binary 3D model"),
        ("WebXR API", "AR placement"),
        ("iOS Quick Look", "iPhone AR"),
        ("Scene Viewer", "Android AR"),
        ("WebGL fallback", "Тогтворгүй үед 3D rotate"),
    ]
    ty = Inches(2.15)
    for k, v in tech_items:
        add_round_rect(s, rx, ty, Inches(4.5), Inches(0.55),
                       LIGHT, line_color=ORANGE)
        add_text(s, rx + Inches(0.15), ty, Inches(2.2), Inches(0.55),
                 k, size=11, bold=True, color=ORANGE, font="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(2.4), ty, Inches(2.0), Inches(0.55),
                 v, size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        ty += Inches(0.65)

    footer(s, 12)


# ════════════════════════════════════════════════════════════════
#  SLIDE 11b — AR жинхэнэ ажиллагаа (2 screenshot)
# ════════════════════════════════════════════════════════════════
def slide_11b_ar_screens():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "AR — жинхэнэ ажиллагаа",
           "Хоолны 3D загварыг бодит ширээний дээр placement хийсэн байдал")

    img_w = Inches(5.6); img_h = Inches(4.4)
    px = [Inches(0.7), Inches(7.05)]
    py = Inches(1.7)
    captions = [
        ("ar-food-1.jpeg", "Хоол №1 — AR placement"),
        ("ar-food-2.jpeg", "Хоол №2 — AR placement"),
    ]
    for i, (fname, cap) in enumerate(captions):
        add_picture_fit(s, fig(fname), px[i], py, img_w, img_h,
                        fallback_text=f"[{fname}]")
        add_text(s, px[i], py + img_h + Inches(0.1), img_w, Inches(0.35),
                 cap, size=12, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)

    # Highlight note
    note_y = Inches(6.65)
    add_round_rect(s, Inches(0.5), note_y, Inches(12.3), Inches(0.55),
                   ORANGE)
    add_text(s, Inches(0.7), note_y, Inches(11.9), Inches(0.55),
             "✨  Browser-аас шууд ажиллана — native app шаардлагагүй.  "
             "<model-viewer> + .glb + WebXR/Quick Look/Scene Viewer.",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 13)


# ════════════════════════════════════════════════════════════════
#  SLIDE 12 — AR ажиллах явц (6-алхамт)
# ════════════════════════════════════════════════════════════════
def slide_12_ar_flow():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "AR ажиллах явц — алхам алхмаар",
           "Зочин хоолны 3D загварыг бодит ширээний дээр харах")

    steps = [
        ("📱", "1. Цэс нээх",
         "Зочин QR scan-ийн дараа\nGuest MenuPage нээгдэнэ",
         PRIMARY),
        ("🥢", "2. AR-тай хоол сонгох",
         "model_url-тай хоолны мөрд\n⟨AR view⟩ товч идэвхтэй",
         ORANGE),
        ("🎬", "3. View товч дарах",
         "<model-viewer> компонент\n.glb файл ачаалагдана",
         GREEN),
        ("📷", "4. Камераар хандах",
         "Browser permission асууна\n(iOS Quick Look / Android Scene Viewer)",
         PURPLE),
        ("🍽", "5. Ширээний дээр placement",
         "AR-аар хоолыг шилжүүлж\nбодит хэмжээгээр харна",
         TEAL),
        ("⊕", "6. Захиалга нэмэх",
         "AR view хаагаад\n'⊕ нэмэх' товчоор сагсанд оруулна",
         RED),
    ]
    cols = 3; rows = 2
    cw = Inches(4.05); ch = Inches(2.3)
    sx = Inches(0.5); sy = Inches(1.7)
    gx = Inches(0.15); gy = Inches(0.2)
    for i, (icon, title, body, color) in enumerate(steps):
        c = i % cols; r = i // cols
        x = sx + c * (cw + gx)
        y = sy + r * (ch + gy)
        # Card frame
        add_round_rect(s, x, y, cw, ch, LIGHT, line_color=color, line_pt=1.5)
        # Header band
        add_rect(s, x, y, cw, Inches(0.55), color)
        add_text(s, x + Inches(0.15), y, Inches(0.6), Inches(0.55),
                 icon, size=22, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), y, cw - Inches(0.8), Inches(0.55),
                 title, size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_text(s, x + Inches(0.15), y + Inches(0.65),
                 cw - Inches(0.3), ch - Inches(0.75),
                 body, size=11, color=DARK, line_spacing=1.3)

    # Limitations / note
    note_y = Inches(6.55)
    add_round_rect(s, Inches(0.5), note_y, Inches(12.3), Inches(0.55),
                   GOLD)
    add_text(s, Inches(0.7), note_y, Inches(11.9), Inches(0.55),
             "⚠  Тэмдэглэл: AR placement зөвхөн ARCore / ARKit дэмжсэн "
             "утаснууд дээр ажиллана. Бусдад WebGL 3D rotate fallback хийгдэнэ.",
             size=11, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 14)


# ════════════════════════════════════════════════════════════════
#  SLIDE 13 — Туршилт ба үр дүн
# ════════════════════════════════════════════════════════════════
def slide_13_results():
    s = add_slide()
    set_bg(s, WHITE)
    header(s, "Туршилт ба үр дүн",
           "Демо орчны preliminary үнэлгээ")

    # Top stats
    stats = [
        ("~75%",  "Захиалгын\nхугацаа буурсан", GREEN),
        ("5x",    "Захиалгын алдаа\nцөөрсөн", ORANGE),
        ("84.2",  "SUS оноо\n(100-аас)", PRIMARY),
        ("8.5 мс","WS round-trip\n(loopback)", PURPLE),
    ]
    sw_ = Inches(2.85); sh_ = Inches(1.7); sy = Inches(1.7)
    sxs = [Inches(0.55), Inches(3.55), Inches(6.55), Inches(9.55)]
    for i, (val, lbl, color) in enumerate(stats):
        add_round_rect(s, sxs[i], sy, sw_, sh_, WHITE,
                       line_color=color, line_pt=2)
        add_text(s, sxs[i], sy + Inches(0.15), sw_, Inches(0.9),
                 val, size=36, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sxs[i], sy + Inches(1.1), sw_, sh_ - Inches(1.2),
                 lbl, size=11, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.2)

    # Mid: 8-scenario functional test results
    add_text(s, Inches(0.55), Inches(3.6), Inches(12.3), Inches(0.4),
             "Функциональ туршилт — 8 хэрэглээний скрипт:",
             size=13, bold=True, color=PRIMARY)
    make_table(s, Inches(0.55), Inches(4.05), Inches(12.3), Inches(2.6),
               ["№", "Хэрэглээний скрипт", "Үр дүн", "Тэмдэглэл"],
               [
                   ["1", "Менежер → Ширээ үүсгэх → QR хэвлэх",        "✓ Pass", "QR PDF гарсан"],
                   ["2", "Утсаар QR scan → Нэр → Цэс үзэх",            "✓ Pass", "Камераар"],
                   ["3", "3 хоол сонгож → Захиалга илгээх",            "✓ Pass", "Bottom sheet сагс"],
                   ["4", "Chef-д real-time захиалга харагдсан эсэх",    "✓ Pass", "Socket.IO ~150мс"],
                   ["5", "Chef 'Бэлдэх' → 'Бэлэн' status солих",        "✓ Pass", "Badge real-time"],
                   ["6", "Cashier төлбөр баталгаажуулах",              "✓ Pass", "Cash/Card сонгох"],
                   ["7", "Менежер dashboard-аас орлого/тайлан",        "✓ Pass", "Recharts визуал"],
                   ["8", "AR view-тай хоол placement",                  "✓ Pass", "ARCore Android"],
               ],
               header_size=10, body_size=9,
               col_widths=[0.5, 5.5, 1.4, 3.5])

    # Bottom note
    add_text(s, Inches(0.55), Inches(6.8), Inches(12.3), Inches(0.35),
             "Тэмдэглэл: эдгээр нь демо орчны preliminary хэмжээс. "
             "Бодит ресторанд явуулах field study нь ирээдүйн ажил.",
             size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s, 15)


# ════════════════════════════════════════════════════════════════
#  SLIDE 14 — Дүгнэлт + Q&A
# ════════════════════════════════════════════════════════════════
def slide_14_thanks():
    s = add_slide()
    set_bg(s, PRIMARY)
    # Big text
    add_text(s, Inches(0), Inches(1.0), SW, Inches(0.7),
             "✓ Дипломын ажлын дүгнэлт", size=20, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Achievements quick
    add_round_rect(s, Inches(1.5), Inches(1.85), Inches(10.3), Inches(2.4),
                   LIGHT, line_color=GOLD)
    achievements = [
        ("✓", "43 REST endpoint + 12 Socket.IO event + 14 PostgreSQL хүснэгт хэрэгжсэн"),
        ("✓", "Multi-guest QR session + AR 3D viewer + real-time notification бүхий бүрэн систем"),
        ("✓", "Демо туршилтад захиалгын хугацаа ~75% буурч, алдаа 5x цөөрсөн"),
        ("✓", "Open-source, монгол хэлтэй, локал хостингтой — Монголын зах зээлд бэлэн"),
    ]
    ay = Inches(2.0)
    for mark, txt in achievements:
        add_text(s, Inches(1.7), ay, Inches(0.4), Inches(0.5),
                 mark, size=18, bold=True, color=GREEN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.2), ay, Inches(9.5), Inches(0.55),
                 txt, size=13, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        ay += Inches(0.55)

    add_text(s, Inches(0), Inches(4.5), SW, Inches(0.5),
             "АНХААРАЛ ТАВЬСАНД",
             size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0), Inches(5.0), SW, Inches(0.8),
             "БАЯРЛАЛАА",
             size=56, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0), Inches(6.0), SW, Inches(0.4),
             "Асуулт байвал хариулахад бэлэн байна",
             size=14, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, 0, SH - Inches(0.5), SW, Inches(0.04), GOLD)
    add_text(s, Inches(0), SH - Inches(0.45), SW, Inches(0.4),
             "Эрдэнэтөгсийн Очбадрах   ·   2026 он   ·   Өмнөговь ТДС",
             size=11, italic=True, color=LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────── Speaker notes (хэлэх үг) ───────────────────────
# Хамгаалалт 10-15 мин. Слайд бүрд тэмдэглэл нэмж, "[ДЭЛГЭРЭНГҮЙ]" болон
# "[ДУРДАЖ ӨНГӨРӨХ]" гэсэн тэмдэглэгээгээр тэргүүлэх ач холбогдлыг ялгасан.
SPEAKER_NOTES = {
1: """[ДУРДАЖ ӨНГӨРӨХ — 30 сек]

Эрхэм багш нар, комиссын гишүүд ээ. Миний бие Эрдэнэтөгсийн Очбадрах бакалаврын дипломын ажлаа танилцуулъя. Сэдэв нь "QR код дээр суурилсан ресторан захиалгын систем — Хан Гарид". Удирдагч багш П. Нарангэрэл. 2026 он.

Цаг хэмнэхийн тулд илтгэлээ системийн ажиллах урсгал дээр гол анхаарлаа төвлөрүүлэн ярих болно.""",

2: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 90 сек]

Ресторан зочлох үед бид бүгд эдгээр асуудалтай тулгардаг:
• Зөөгчийг удаан хүлээх
• Цэс нөөшилж байх үед сольсон гэдгийг мэдэхгүй
• Төлбөр хариуцах нэг зөөгч завгүй байх
• Гадаад зочинд орчуулга хийх хүндрэл
• COVID-ийн дараах эрүүл ахуйн санаа зовнил

Дэлхийн зах зээлд QR захиалгын систем нь үйлчилгээний хугацааг 25%-иар бууруулдаг, дундаж захиалгын дүнг 18%-иар нэмэгдүүлдэг гэж ResearchGate болон McKinsey-ийн судалгаанд тогтоосон.

Монгол улсад 2024 онд хүнсний үйлчилгээний салбар ДНБ-ний 4.2% эзэлсэн ч QR захиалга нэвтрүүлсэн ресторан 5%-аас бага. Зах зээл нээлттэй, эрэлт өсч байна.""",

3: """[ДУРДАЖ ӨНГӨРӨХ — 60 сек]

Дэлхийн зах зээлд Toast, Lightspeed, Square POS зэрэг шийдэл байна. Гэхдээ:
• Сард 80+ доллар — Монголын ресторанд үнэтэй
• Зөвхөн англи хэлтэй
• AR functionality байхгүй
• Multi-guest session дэмждэггүй

Манай систем энэ 4 цоорхойг хааж байна: үнэгүй опен-сорс, монгол хэлтэй, AR-тэй, нэг QR-аар олон зочин нэгдэх боломжтой. Хүснэгтэн дээр 8 шалгуураар харьцуулсан байна.""",

4: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 75 сек]

Зорилго: монгол хэлтэй, AR-тэй, real-time, multi-guest QR ресторан захиалгын систем хөгжүүлэх.

Архитектур нь layered загвараар 4 давхаргатай:
• ҮЗҮҮЛБЭР ДАВХАРГА — React 19 + Vite 7 + TanStack Query + Tailwind CSS
• API ДАВХАРГА — Express 5 + Socket.IO 4 + JWT auth
• ӨГӨГДЛИЙН ДАВХАРГА — PostgreSQL 16 + Drizzle ORM, 14 хүснэгт
• ГАДААД ҮЙЛЧИЛГЭЭ — Cloudflare tunnel, model-viewer AR

Бүх давхрага хооронд тодорхой interface. Энэ нь maintainability болон scalability-г хангаж байгаа гол шийдэл.""",

5: """[ДУРДАЖ ӨНГӨРӨХ — 60 сек]

Эдгээр нь жинхэнэ систем дээрх 3 mobile screenshot.

ЗҮҮН — Зочин QR-аа уншуулаад цэс рүү орлоо. Категори chip-ээр шүүх, "AR view" товчоор хоолоо урьдчилан 3D-ээр харах боломжтой.

ДУНД — Сагсанд хоолоо хийгээд тоо ширхэгээ нэмэх/хасах. Нийт дүн доороо real-time-аар тоологдоно.

БАРУУН — Захиалга илгээсний дараа автоматаар очдог "захиалгын төлөв" хуудас. "Бэлдэж буй" → "Бэлэн" гэсэн алхамууд real-time-аар шинэчлэгдэнэ — refresh хийх шаардлагагүй.""",

6: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 60 сек]

ЗҮҮН ТАЛД — захиалгын самбар (/admin). Орж ирсэн захиалгуудыг Kanban загвараар "Бэлдэж буй / Бэлэн / Хүргэгдсэн" гэсэн 3 баганаар хувааж байгаа. Гал тогооч төлвийг дарахад зочны утсанд real-time-аар badge харагдана.

БАРУУН ТАЛД — ширээний удирдлага. Менежер шинэ ширээ нэмэхэд систем UUID-аар qr_token автоматаар үүсгэж, QR кодыг шууд хэвлэх боломжтой PNG болгож өгнө. Энэ QR-уудыг хэвлээд ширээн дээр тогтооно.""",

7: """[ДУРДАЖ ӨНГӨРӨХ — 45 сек]

Энд менежерийн нэмэлт 3 модулийг харуулж байна.

ЗҮҮН — Ажилтны нэвтрэх дэлгэц. JWT баталгаажуулалт, bcrypt password hash. Manager / Chef / Cashier гэсэн 3 дүртэй, role-based middleware.

ДУНД — Цэс удирдлага. Хоол нэмэх, зураг ачаалах, .glb AR model хавсаргах боломжтой.

БАРУУН — Тайлангийн дэлгэц. Өдрийн орлого, эрэлттэй хоолны жагсаалт, цаг тутмын ачаалал — бүгд Recharts library-ээр визуал болгож үзүүлж байна.""",

8: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — энэ слайд хамгийн чухал — 100 сек]

Зочны бүрэн ажиллах урсгал 7 алхамтай:

1. QR scan — Зочин ширээнийхээ QR-ыг утсаар уншина
2. Session үүсэх — Сервер qr_token-ийг шалгаж JWT олгоно
3. Цэс — Хоол сонгож сагсанд оруулна, хүсвэл AR-аар үзнэ
4. Илгээх — POST /api/cart дуудлага, DB-д захиалга үүсэх, inventory автомат хасагдах
5. Real-time event — Socket.IO order:new event бүх ажилтан руу broadcast
6. Бэлдэлт — Chef Kanban-аас 'Бэлэн' даравал зочинд утсанд push мэдэгдэл
7. Төлбөр — Cashier QPay / бэлэн / карт сонгож, session хаагдах, ширээ available болох

ГОЛ ОНЦЛОГ: 5 ба 6-р алхамд ХҮНИЙ ОРОЛЦОО БАЙХГҮЙ. Захиалга, төлөв шинэчлэлт, мэдэгдэл бүгд автомат, real-time. Энэ нь ажилтны ачааллыг хамгийн их буулгасан хэсэг.""",

9: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 90 сек]

Энэ бол манай системийн ХАМГИЙН ОНЦЛОГ ДАВУУ ТАЛ.

Уламжлалт системд нэг хүн нэг утсаар захиалдаг. Бидний системд 3, 5, 10 хүн нэг QR-аар нэгдэж тус тусын утаснаасаа сонгож, ХАМТЫН САГСАНД оруулж захиалдаг.

Техникийн талаас 3 PostgreSQL хүснэгт хамтран ажиллана:
• table_sessions — ширээний идэвхтэй session
• session_participants — хэн хэн нь нэгдсэн
• shared_cart_items — хамтарсан сагс

Real-time-аар нэг хүн item нэмэхэд бусдын утсанд шууд харагдана — Socket.IO session:<id> room ашиглаж.

Use case: гэр бүлийн оройн хоол, ажилчдын lunch, найзуудтай уулзалт — group ordering-д төгс тохирно. Toast, Square зэрэг гадаад систем энэ функцийг дэмждэггүй.""",

10: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 60 сек]

Socket.IO 4-р хувилбараар real-time event дамжуулна. Гурван логик room:
• "staff" — бүх ажилтан
• "chef" — зөвхөн гал тогооч
• "session:<id>" — тухайн ширээний зочид

Sequence diagram-аас харахад: зочин 'Захиалга илгээх' дарахад POST /api/cart явна → DB write → Socket.IO order:new event → chef room руу broadcast. Refresh шаардлагагүй.

Хэмжсэн дундаж latency: loopback ~8.5 мс, Cloudflare tunnel-ээр ~142 мс, P95 195 мс. Энэ нь хүний нүдэнд "шууд" мэдрэгдэх түвшин (200 мс-ээс доош).""",

11: """[ДУРДАЖ ӨНГӨРӨХ — 45 сек]

Менежерийн самбарт 9 үндсэн модуль:
1. Цэс удирдах
2. Ширээ + QR хэвлэх
3. Нөөц хяналт — доод түвшинд алерт
4. Ажилтан удирдлага
5. Орлогын тайлан
6. Банер карусел
7. Зочдын сэтгэгдэл
8. Захиалгын түүх
9. Тохиргоо

Эдгээр нь CRUD endpoint + Role-based middleware-ээр хамгаалагдсан.""",

12: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 60 сек]

AR ашиглах болсон шалтгаан: зочин хоолоо захиалахаас өмнө хэмжээ, өнгө, бүтцийг 3D-ээр харах. Google-ийн e-commerce судалгаагаар AR ашигласан үед буцаалт 40%-иар буурдаг, conversion 94%-иар өсдөг.

Технологийн стек:
• <model-viewer> — Google-ийн нээлттэй web component
• .glb формат — GLTF binary 3D model
• WebXR API — placement
• iOS Quick Look + Android Scene Viewer — нэмэлт config-гүй шууд ажиллана
• WebGL fallback — тогтворгүй утсанд 3D rotate

Энэ нь plug-and-play, native app шаардахгүй — browser-ээс шууд.""",

13: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 60 сек]

Энэ бол AR-ийн жинхэнэ ажиллагааны 2 screenshot. Зүүн талд нэг хоол, баруун талд өөр хоол — хоёул зочны ширээн дээр placement хийгдсэн байдал.

Бодит хэмжээгээр харагдаж байгаагаа анхаарна уу — Blender дээр boundary box-ыг empirical-ээр 29% (scale 0.000435) болгож тааруулсан учир хоолны жинхэнэ диаметртэй (10-25 см) ойролцоо.

Зочин AR view-ийг камераар нээгээд гар хөдөлгөж placement хийнэ. Хоолны өнгө, гарнитур, тавиурын өндрийг бодитойгоор үзэх боломж. Гадаад зочдод нэр ойлгомжгүй ч зургаар сонгох шинэ боломж.

Native app огт шаардлагагүй — browser-ээс шууд. iPhone дээр Safari Quick Look, Android дээр Chrome Scene Viewer автоматаар идэвхждэг.""",

14: """[ДУРДАЖ ӨНГӨРӨХ — 45 сек]

AR ажиллах 6 алхам:
1. Зочин цэс нээнэ
2. AR view идэвхтэй хоолыг сонгоно (model_url тохируулсан хоолны мөрд товч идэвхтэй)
3. View товч даравал .glb ачаалагдана
4. Browser камерын зөвшөөрөл асууна
5. Камераар ширээний дээр placement, бодит хэмжээгээр харагдана
6. AR view хааж захиалга нэмнэ

iOS дээр Quick Look, Android дээр Scene Viewer автоматаар нээгдэнэ.""",

15: """[ДЭЛГЭРЭНГҮЙ ЯРИХ — 90 сек]

Хэмжсэн 4 гол үзүүлэлт:

1. ФУНКЦИОНАЛ ТУРШИЛТ — 8 test script, бүгд PASS:
   QR session нээх, цэс үзэх, сагс, multi-guest нэгдэх, Socket realtime, AR view, тайлан, нэвтрэлт

2. ХУГАЦАА — захиалгын дундаж 4 минут → 1 минут (75% бууралт)

3. АЛДАА — зөөгчийн алдаа 5 удаа/өдөр → 1 удаа/өдөр (5 дахин цөөрсөн)

4. SUS ОНОО — System Usability Scale 82/100. Industry дунж 68. "Good"-ээс "Excellent" руу шилжсэн.

Анхааруулга: эдгээр нь демо орчны preliminary хэмжээс. Бодит ресторанд field study хийх нь ирээдүйн ажил байна.""",

16: """[ДУРДАЖ ӨНГӨРӨХ — хаах үг, 30 сек]

Дүгнэлт:
• 43 REST endpoint, 12 Socket.IO event, 14 PostgreSQL хүснэгт хэрэгжсэн
• Multi-guest, AR, real-time бүх онцлог ажиллаж байна
• Захиалгын хугацаа 75% буурч, алдаа 5 дахин цөөрсөн
• Open-source, монгол хэлтэй, локал хостингтой — Монголын зах зээлд бэлэн

Анхаарал тавьсанд баярлалаа.

Асуултад хариулахад бэлэн байна.

ХҮЛЭЭГДЭЖ БУЙ АСУУЛТУУД (бэлдсэн):
• Хэрэв хэрэглэгч ачаалал нэмэгдвэл? → horizontal scale, Redis adapter
• Аюулгүй байдал? → JWT + bcrypt + RBAC + CORS + rate limit
• Bug бий юу? → 8 test PASS, гэхдээ field study хэрэгтэй
• Үнэ? → Open-source, hosting сард 5-15 доллар
• AR хувилбаргүй утсанд? → WebGL 3D rotate fallback""",
}


def add_speaker_notes():
    """Илтгэлд хэлэх текстийг слайд бүрд speaker note болгож суулгана."""
    for idx, slide in enumerate(prs.slides, start=1):
        text = SPEAKER_NOTES.get(idx)
        if text is None:
            continue
        slide.notes_slide.notes_text_frame.text = text


# ─────────────────────── Бүх слайдыг угсрах ───────────────────────
def build():
    slide_1_title()
    slide_2_rationale()
    slide_3_comparison()
    slide_4_goal_arch()
    slide_5_screens_public()
    slide_6_screens_staff()
    slide_6b_screens_admin_more()
    slide_7_workflow()
    slide_8_multiguest()
    slide_9_realtime()
    slide_10_manager()
    slide_11_ar_benefits()
    slide_11b_ar_screens()
    slide_12_ar_flow()
    slide_13_results()
    slide_14_thanks()

    add_speaker_notes()

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "workflow.pptx")
    prs.save(out)
    try:
        print(f"[OK] PPT amjilttai uusgegdlee: {out}")
        print(f"     Slide count: {TOTAL_SLIDES}")
        print(f"     Speaker notes: {len(SPEAKER_NOTES)} slide(s)")
    except UnicodeEncodeError:
        pass
    return out


if __name__ == "__main__":
    build()
