"""Generate a 20-slide academic presentation describing the Restaurant Table Booking system."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "Figures"
OUT = ROOT / "Restaurant-Booking-Presentation.pptx"

# Color palette
PRIMARY = RGBColor(0x18, 0x4A, 0xC4)        # blue
ACCENT = RGBColor(0xE6, 0x85, 0x10)          # amber
DARK = RGBColor(0x1F, 0x2D, 0x3D)
LIGHT = RGBColor(0xF6, 0xF8, 0xFC)
MUTED = RGBColor(0x6B, 0x76, 0x85)
SUCCESS = RGBColor(0x10, 0xA3, 0x6E)
DANGER = RGBColor(0xD3, 0x36, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PHONE_BODY = RGBColor(0x22, 0x33, 0x49)
PHONE_BORDER = RGBColor(0x10, 0x18, 0x28)

FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ============================================================
# Helpers
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_round(slide, x, y, w, h, fill, line=None, adj=0.12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = adj
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK,
                bullet_color=PRIMARY, space_after=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        bullet = p.add_run()
        bullet.text = "▸ "
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        bullet.font.name = FONT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.15), SW - Inches(1), Inches(0.6),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.5), SW - Inches(1), Inches(0.4),
                 subtitle, size=12, color=RGBColor(0xE0, 0xE9, 0xF7))


def footer(slide, num):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), LIGHT)
    add_text(slide, Inches(0.4), SH - Inches(0.35), Inches(8), Inches(0.35),
             "Рестораны ширээ захиалга, цэс удирдлагын мобайл систем", size=10,
             color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(1), SH - Inches(0.35), Inches(0.6), Inches(0.35),
             f"{num} / 20", size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT)


def add_phone_frame(slide, image_path, x, y, height, label=None):
    """Draw a smartphone frame around the screenshot.
    Returns (body_x, body_y, body_w, body_h)."""
    aspect = 0.49  # width / height (modern phone with frame)
    width = height * aspect

    # Outer body (dark)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    body.adjustments[0] = 0.09
    body.fill.solid()
    body.fill.fore_color.rgb = PHONE_BODY
    body.line.color.rgb = PHONE_BORDER
    body.line.width = Pt(0.5)
    body.shadow.inherit = False

    # Side button (right)
    btn_h = height * 0.08
    btn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + width - Emu(int(Pt(0.5).emu * 2)),
                                 y + height * 0.18,
                                 Inches(0.04), btn_h)
    btn.fill.solid()
    btn.fill.fore_color.rgb = PHONE_BORDER
    btn.line.fill.background()
    btn.shadow.inherit = False

    # Screen area inset
    inset_xr = 0.03
    inset_yr = 0.018
    screen_x = x + width * inset_xr
    screen_y = y + height * inset_yr
    screen_w = width * (1 - 2 * inset_xr)
    screen_h = height * (1 - 2 * inset_yr)

    # Screen background
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    screen_x, screen_y, screen_w, screen_h)
    screen.adjustments[0] = 0.06
    screen.fill.solid()
    screen.fill.fore_color.rgb = WHITE
    screen.line.fill.background()
    screen.shadow.inherit = False

    # Image inside screen
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), screen_x, screen_y,
                                 width=screen_w, height=screen_h)

    # Notch (speaker pill at top)
    notch_w = width * 0.30
    notch_h = height * 0.018
    notch_x = x + (width - notch_w) / 2
    notch_y = y + height * 0.008
    notch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   notch_x, notch_y, notch_w, notch_h)
    notch.adjustments[0] = 0.5
    notch.fill.solid()
    notch.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x18)
    notch.line.fill.background()
    notch.shadow.inherit = False

    # Optional caption below
    if label:
        add_text(slide, x - Inches(0.5), y + height + Inches(0.1),
                 width + Inches(1.0), Inches(0.4),
                 label, size=11, bold=True, color=MUTED,
                 align=PP_ALIGN.CENTER)

    return (x, y, width, height)


def add_image(slide, path, x, y, w=None, h=None):
    if path.exists():
        if w and h:
            return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        if w:
            return slide.shapes.add_picture(str(path), x, y, width=w)
        if h:
            return slide.shapes.add_picture(str(path), x, y, height=h)
        return slide.shapes.add_picture(str(path), x, y)
    box = add_round(slide, x, y, w or Inches(4), h or Inches(3), LIGHT, MUTED)
    add_text(slide, x, y, w or Inches(4), h or Inches(3),
             f"[{path.name}]", size=14, color=MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return box


# ============================================================
# SLIDE 1 — Гарчиг
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PRIMARY)
add_oval(s, Inches(-2), Inches(-2), Inches(6), Inches(6),
         RGBColor(0x32, 0x66, 0xD8))
add_oval(s, SW - Inches(4), SH - Inches(4), Inches(6), Inches(6),
         RGBColor(0x32, 0x66, 0xD8))

add_text(s, Inches(0.8), Inches(0.6), SW - Inches(1.6), Inches(0.5),
         "БОЛОВСРОЛЫН ИХ СУРГУУЛЬ", size=12, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(0.5),
         "БАКАЛАВРЫН ДИПЛОМЫН АЖИЛ", size=14, color=WHITE, bold=True)

add_rect(s, Inches(0.8), Inches(1.7), Inches(0.8), Inches(0.05), ACCENT)

add_text(s, Inches(0.8), Inches(2.0), SW - Inches(1.6), Inches(2.5),
         "Рестораны ширээ захиалга,\nцэс удирдлагын мобайл\nсистемийн загварчлал, хэрэгжилт",
         size=36, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
         "QR кодод суурилсан, бодит цагийн харилцаатай, олон төхөөрөмж дэмжсэн систем",
         size=16, color=RGBColor(0xE0, 0xE9, 0xF7))

add_text(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.4),
         "Дипломант:", size=12, color=RGBColor(0xC0, 0xD0, 0xEC))
add_text(s, Inches(2.5), Inches(5.7), Inches(11), Inches(0.4),
         "Ц. Эрхэмбаяр", size=14, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4),
         "Удирдагч багш:", size=12, color=RGBColor(0xC0, 0xD0, 0xEC))
add_text(s, Inches(2.5), Inches(6.1), Inches(11), Inches(0.4),
         "............................", size=14, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
         "Улаанбаатар хот · 2026 он", size=12, color=RGBColor(0xC0, 0xD0, 0xEC))


# ============================================================
# SLIDE 2 — Судалгааны үндэслэл
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "1. Судалгааны ажлын үндэслэл")

add_text(s, Inches(0.6), Inches(1.1), SW - Inches(1.2), Inches(1.0),
         "Орчин үеийн ресторанууд хэрэглэгчийн үйлчилгээ, удирдлагын үр ашиг, "
         "өгөгдлийн нэгтгэл зэрэг олон тулгамдсан асуудалтай тулгардаг. "
         "Дараах бэрхшээлүүдийг шийдвэрлэх шаардлага үүссэн:",
         size=13, color=DARK)

problems = [
    ("Үйлчилгээний хугацаа", DANGER,
     "Уламжлалт цаасан цэс, гар бичмэл захиалгын процесст 4–6 минут зарцуулагдаж, "
     "захиалга бүрд үйлчлэгчийн зайлшгүй оролцоо шаардагддаг."),
    ("Захиалгын алдаа", DANGER,
     "Гар ажиллагаатай захиалга нь буруу хүлээн авах, давхардах, алга болох эрсдэлтэй "
     "(дунджаар 6–10%-ийн алдаа гардаг)."),
    ("Удирдлагын мэдээллийн хязгаарлалт", DANGER,
     "Менежер өдрийн орлого, эрэлттэй хоол, нөөцийн тоог гар аргаар нэгтгэх ёстой "
     "болж шийдвэр гаргалт удааширдаг."),
    ("Зочны хязгаарлагдмал хандалт", DANGER,
     "Зочин ресторанд биеэр ирэхээс өмнө цэс үзэх, ширээ урьдчилан захиалах, "
     "сэтгэгдэл уншихад хүндрэлтэй."),
]
y = Inches(2.3)
for t, color, d in problems:
    add_round(s, Inches(0.6), y, SW - Inches(1.2), Inches(1.05), LIGHT)
    add_rect(s, Inches(0.6), y, Inches(0.15), Inches(1.05), color)
    add_text(s, Inches(0.95), y + Inches(0.1),
             SW - Inches(1.5), Inches(0.4),
             t, size=15, bold=True, color=DARK)
    add_text(s, Inches(0.95), y + Inches(0.45),
             SW - Inches(1.5), Inches(0.6),
             d, size=11, color=MUTED)
    y += Inches(1.15)

footer(s, 2)


# ============================================================
# SLIDE 3 — Зорилго ба зорилт
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "2. Судалгааны зорилго ба зорилт")

# Goal box
add_round(s, Inches(0.6), Inches(1.2), SW - Inches(1.2), Inches(1.4), LIGHT)
add_rect(s, Inches(0.6), Inches(1.2), Inches(0.2), Inches(1.4), PRIMARY)
add_text(s, Inches(1.0), Inches(1.3), SW - Inches(1.6), Inches(0.5),
         "Судалгааны зорилго", size=16, bold=True, color=PRIMARY)
add_text(s, Inches(1.0), Inches(1.8), SW - Inches(1.6), Inches(0.9),
         "QR кодод суурилсан, бодит цагийн харилцаатай, олон төхөөрөмж дэмждэг "
         "рестораны ширээ захиалга, цэс удирдлагын мобайл системийг загварчилж, "
         "хэрэгжүүлэх замаар үйлчилгээний чанарыг сайжруулах, ажиллагааны үр ашгийг "
         "нэмэгдүүлэх.",
         size=12, color=DARK)

# Objectives
add_text(s, Inches(0.6), Inches(2.9), SW - Inches(1.2), Inches(0.5),
         "Зорилтууд", size=18, bold=True, color=PRIMARY)

objectives = [
    ("1", "Онолын судалгаа",
     "QR код, WebSocket, REST API, JWT, ORM зэрэг технологийн онолын суурийг судлах."),
    ("2", "Шаардлагын анализ",
     "Хэрэглэгчийн үүрэг, функционал болон функционал бус шаардлагыг тодорхойлох."),
    ("3", "Системийн загвар боловсруулах",
     "ER, UML, дарааллын диаграм бүхий бүрэн архитектурын загвар гаргах."),
    ("4", "Систем хэрэгжүүлэх",
     "4 давхаргат архитектурын дагуу REST API, мобайл клиент, удирдлагын самбар "
     "хөгжүүлэх."),
    ("5", "Туршилт ба үнэлгээ",
     "k6 ачааллын туршилт, SUS хэрэглэгчийн үнэлгээгээр системийн чанарыг шалгах."),
]
oy = Inches(3.5)
for num, title, desc in objectives:
    add_round(s, Inches(0.6), oy, SW - Inches(1.2), Inches(0.65), LIGHT)
    add_oval(s, Inches(0.7), oy + Inches(0.07),
             Inches(0.5), Inches(0.5), PRIMARY)
    add_text(s, Inches(0.7), oy + Inches(0.07), Inches(0.5), Inches(0.5),
             num, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.4), oy + Inches(0.05),
             Inches(3.5), Inches(0.55),
             title, size=13, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.0), oy + Inches(0.05),
             SW - Inches(5.6), Inches(0.55),
             desc, size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    oy += Inches(0.7)

footer(s, 3)


# ============================================================
# SLIDE 4 — Ижил төстэй системийн харьцуулсан судалгаа
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "3. Ижил төстэй системийн харьцуулсан судалгаа")

# Comparison table
table_data = [
    ["Шинж чанар", "Toast POS", "Lightspeed", "Square POS",
     "Ajisen Ramen", "Бидний систем"],
    ["QR кодоор захиалга", "Хязгаарлагдмал", "Үгүй", "Үгүй", "Тийм", "Тийм"],
    ["Олон төхөөрөмжийн хамтын сагс", "Үгүй", "Үгүй", "Үгүй", "Үгүй", "Тийм"],
    ["Бодит цагийн харилцаа", "Хязгаарлагдмал", "Үгүй", "Үгүй", "Хязгаарт.", "Тийм"],
    ["Public хандалтын хуудас", "Үгүй", "Үгүй", "Үгүй", "Хязгаарт.", "Тийм"],
    ["Монгол хэлний дэмжлэг", "Үгүй", "Үгүй", "Үгүй", "Үгүй", "Тийм"],
    ["Сарын зардал (USD)", "165+", "89+", "60+", "—", "10–25"],
    ["Нээлттэй эх", "Үгүй", "Үгүй", "Үгүй", "Үгүй", "Тийм"],
]

tx = Inches(0.4)
ty = Inches(1.15)
col_widths = [Inches(3.4), Inches(1.6), Inches(1.6), Inches(1.6),
              Inches(1.6), Inches(2.3)]
row_h = Inches(0.45)

for r, row in enumerate(table_data):
    cx = tx
    is_header = r == 0
    for c, val in enumerate(row):
        if is_header:
            bg = PRIMARY
            fg = WHITE
        elif c == 5:
            bg = RGBColor(0xDC, 0xF2, 0xE5)
            fg = SUCCESS
        else:
            bg = LIGHT if r % 2 == 0 else WHITE
            fg = DARK
        add_rect(s, cx, ty + r * row_h, col_widths[c], row_h, bg,
                 RGBColor(0xCC, 0xCC, 0xCC))
        size = 11 if is_header else 10
        bold = is_header or c == 0 or c == 5
        align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        add_text(s, cx + Inches(0.1), ty + r * row_h,
                 col_widths[c] - Inches(0.2), row_h,
                 val, size=size, bold=bold, color=fg,
                 anchor=MSO_ANCHOR.MIDDLE, align=align)
        cx += col_widths[c]

# Advantages summary
ay = ty + len(table_data) * row_h + Inches(0.2)
add_round(s, Inches(0.4), ay, SW - Inches(0.8), Inches(2.3),
          RGBColor(0xDC, 0xF2, 0xE5))
add_text(s, Inches(0.6), ay + Inches(0.1), SW - Inches(1.2), Inches(0.5),
         "Боловсруулсан системийн ялгарах давуу талууд", size=14, bold=True,
         color=SUCCESS)
advs = [
    ("Хямд өртөг — гадаадын системээс 4–10 дахин бага сарын зардалтай",
     "Монгол хэлний бүрэн дэмжлэг — интерфейс, мэдэгдэл бүгд монголоор"),
    ("Олон төхөөрөмжийн хамтын сагс — нэг ширээний 4 зочин нэгэн зэрэг",
     "QR-гүй public хуудас — цэс, ширээний захиалга, сэтгэгдэл интернетээс шууд"),
    ("Stable QR (хэвлэх шаардлагагүй) + ephemeral session JWT хослуулсан",
     "Нээлттэй эх код — рестораны өөрийн шаардлагад нийцүүлж тохируулах"),
]
for i, (a, b) in enumerate(advs):
    add_text(s, Inches(0.7), ay + Inches(0.55) + i * Inches(0.45),
             Inches(6.2), Inches(0.4),
             "✓ " + a, size=10, color=DARK)
    add_text(s, Inches(7.0), ay + Inches(0.55) + i * Inches(0.45),
             Inches(6.0), Inches(0.4),
             "✓ " + b, size=10, color=DARK)

footer(s, 4)


# ============================================================
# SLIDE 5 — Системийн архитектур
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "4. Системийн архитектур", "4 давхаргат клиент–серверийн загвар")

layers = [
    ("1. Харуулах давхарга (Presentation Layer)",
     "React 18 · TypeScript · Tailwind CSS v4 · Vite",
     RGBColor(0xD9, 0xE6, 0xFA)),
    ("2. Хэрэглээний давхарга (Application Layer)",
     "Node.js · Express.js · Socket.IO v4 · 43 REST endpoint",
     RGBColor(0xD0, 0xF0, 0xDC)),
    ("3. Бизнес логикийн давхарга (Business Logic Layer)",
     "JWT HS256 · Role-based access · Drizzle ORM · Zod validation",
     RGBColor(0xFD, 0xEB, 0xC9)),
    ("4. Өгөгдлийн давхарга (Data Layer)",
     "PostgreSQL 16 · Drizzle migration · 11 хүснэгт",
     RGBColor(0xFA, 0xD7, 0xD7)),
]
y = Inches(1.3)
for title, tech, color in layers:
    add_round(s, Inches(1.5), y, SW - Inches(3), Inches(1.05), color)
    add_text(s, Inches(1.8), y + Inches(0.1), SW - Inches(3.6), Inches(0.45),
             title, size=17, bold=True, color=DARK)
    add_text(s, Inches(1.8), y + Inches(0.55), SW - Inches(3.6), Inches(0.45),
             tech, size=12, color=MUTED)
    y += Inches(1.2)

# arrows
arr_x = Inches(0.7)
add_text(s, arr_x, Inches(1.4), Inches(0.7), Inches(4.5),
         "↓\n↓\n↓\n↓", size=24, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, SW - Inches(1.2), Inches(1.4), Inches(0.7), Inches(4.5),
         "↑\n↑\n↑\n↑", size=24, color=PRIMARY, align=PP_ALIGN.CENTER)

# Side labels
add_text(s, Inches(0.4), Inches(1.5), Inches(0.6), Inches(0.4),
         "Хүсэлт", size=10, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, SW - Inches(1.0), Inches(1.5), Inches(0.6), Inches(0.4),
         "Хариу", size=10, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

footer(s, 5)


# ============================================================
# SLIDE 6 — Технологийн стек
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "5. Хэрэгжүүлэлтэд ашигласан технологиуд")

stacks = [
    ("Frontend (Клиент тал)", PRIMARY, [
        "React 18 + TypeScript",
        "Vite (dev server, HMR)",
        "Tailwind CSS v4",
        "TanStack Query (server state)",
        "Zustand (client state)",
        "Wouter (routing)",
        "shadcn/ui (компонент сан)",
    ]),
    ("Backend (Сервер тал)", SUCCESS, [
        "Node.js 22 + Express.js",
        "Socket.IO v4 (WebSocket)",
        "JWT HS256 (нэвтрэлт)",
        "Drizzle ORM",
        "Zod (баталгаажуулалт)",
        "bcrypt (нууц үг хеш)",
        "Multer (файл upload)",
    ]),
    ("Infrastructure (Дэд бүтэц)", ACCENT, [
        "PostgreSQL 16",
        "pnpm workspace монорепо",
        "Docker контейнер",
        "Cloudflare Tunnel (демо)",
        "drizzle-kit (migration)",
        "k6 (ачааллын туршилт)",
        "Git (хувилбарын удирдлага)",
    ]),
]
cw = Inches(4.0)
gap = Inches(0.25)
start_x = (SW - (cw * 3 + gap * 2)) / 2
y = Inches(1.3)
for i, (label, color, items) in enumerate(stacks):
    x = start_x + i * (cw + gap)
    add_round(s, x, y, cw, Inches(5.5), LIGHT)
    add_rect(s, x, y, cw, Inches(0.6), color)
    add_text(s, x, y, cw, Inches(0.6),
             label, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.3), y + Inches(0.8), cw - Inches(0.6),
                Inches(4.5), items, size=12, bullet_color=color)

footer(s, 6)


# ============================================================
# SLIDE 7 — Хэрэглэгчдийн төрлүүд
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "6. Хэрэглэгчдийн үүрэг (Roles)")

roles = [
    ("👤", "Зочин (Guest)", PRIMARY,
     ["QR-аар цэс үзэх", "Сагсанд хоол нэмэх", "Захиалга өгөх",
      "Захиалгын статус хянах", "Сэтгэгдэл үлдээх"]),
    ("🧑‍💼", "Кассир (Cashier)", SUCCESS,
     ["Захиалга баталгаажуулах", "Статус шинэчлэх", "Тооцоо хийх",
      "Ширээний QR хэвлэх", "Захиалгын түүх харах"]),
    ("👨‍🍳", "Гал тогоо (Kitchen)", ACCENT,
     ["Бэлдэх дараалал харах", "ready болгож тэмдэглэх",
      "Нөөцийн доод хэмжээний мэдэгдэл хүлээн авах"]),
    ("📊", "Менежер (Manager)", RGBColor(0x6F, 0x42, 0xC1),
     ["Цэс/ширээ удирдлага", "Ажилтан удирдах",
      "Тайлан, орлого", "Сэтгэгдлийн хариулт",
      "Нөөц, банер, тохиргоо"]),
]
cw, ch = Inches(5.8), Inches(2.5)
gap = Inches(0.3)
start_x = (SW - (cw * 2 + gap)) / 2
y = Inches(1.2)
for i, (icon, name, color, perms) in enumerate(roles):
    col = i % 2
    row = i // 2
    x = start_x + col * (cw + gap)
    yy = y + row * (ch + gap)
    add_round(s, x, yy, cw, ch, LIGHT)
    add_rect(s, x, yy, Inches(0.2), ch, color)
    add_text(s, x + Inches(0.4), yy + Inches(0.1),
             Inches(0.8), Inches(0.8),
             icon, size=32)
    add_text(s, x + Inches(1.3), yy + Inches(0.2),
             cw - Inches(1.5), Inches(0.5),
             name, size=17, bold=True, color=color)
    add_bullets(s, x + Inches(1.3), yy + Inches(0.7),
                cw - Inches(1.5), ch - Inches(0.8),
                perms, size=11, bullet_color=color)

footer(s, 7)


# ============================================================
# SLIDE 8 — Public Landing
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "7. Public хандалтын хуудас — Нүүр",
       "QR код шаардлагагүй, интернетээс шууд хандах хуудаснууд")

# Mock browser frame on left
mx, my = Inches(0.5), Inches(1.2)
mw, mh = Inches(7.5), Inches(5.7)
add_round(s, mx, my, mw, mh, LIGHT)
add_rect(s, mx, my, mw, Inches(0.4), RGBColor(0xE0, 0xE0, 0xE0))
for i, c in enumerate([RGBColor(0xFF, 0x5F, 0x57),
                       RGBColor(0xFE, 0xBC, 0x2E),
                       RGBColor(0x28, 0xC8, 0x40)]):
    add_oval(s, mx + Inches(0.1 + i * 0.3),
             my + Inches(0.1), Inches(0.2), Inches(0.2), c)
add_round(s, mx + Inches(1.2), my + Inches(0.07),
          Inches(5.5), Inches(0.26), WHITE)
add_text(s, mx + Inches(1.3), my + Inches(0.05),
         Inches(5.4), Inches(0.3),
         "🔒 restaurant-domain.mn",
         size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

# Hero area
hy = my + Inches(0.6)
add_rect(s, mx + Inches(0.3), hy, mw - Inches(0.6), Inches(2.0), PRIMARY)
add_text(s, mx + Inches(0.5), hy + Inches(0.2),
         mw - Inches(1), Inches(0.4),
         "Тавтай морилно уу 🍽", size=10, color=WHITE)
add_text(s, mx + Inches(0.5), hy + Inches(0.5),
         mw - Inches(1), Inches(0.7),
         "Mongolia Restaurant", size=24, bold=True, color=WHITE)
add_text(s, mx + Inches(0.5), hy + Inches(1.2),
         mw - Inches(1), Inches(0.5),
         "Цэсээ үзэх, ширээ захиалах, сэтгэгдэл үлдээх — нэг дороос",
         size=11, color=WHITE)
add_round(s, mx + Inches(0.5), hy + Inches(1.55),
          Inches(1.7), Inches(0.35), ACCENT)
add_text(s, mx + Inches(0.5), hy + Inches(1.55),
         Inches(1.7), Inches(0.35),
         "Цэс үзэх →", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 3 action cards
cy = hy + Inches(2.3)
labels = [("Цэс үзэх", "🍽"),
          ("Ширээ захиалах", "📅"),
          ("Сэтгэгдэл", "⭐")]
acw = (mw - Inches(0.6) - Inches(0.4)) / 3
for i, (lbl, ic) in enumerate(labels):
    cx = mx + Inches(0.3) + i * (acw + Inches(0.2))
    add_round(s, cx, cy, acw, Inches(1.4), WHITE,
              RGBColor(0xE0, 0xE0, 0xE0))
    add_text(s, cx, cy + Inches(0.15), acw, Inches(0.5),
             ic, size=22, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + Inches(0.7), acw, Inches(0.4),
             lbl, size=12, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)

# Description right
dx, dy, dw = Inches(8.4), Inches(1.5), Inches(4.4)
add_text(s, dx, dy, dw, Inches(0.5), "Гол боломжууд",
         size=18, bold=True, color=PRIMARY)
features = [
    "Hero banner + рестораны нэр (settings API)",
    "3 том CTA: Цэс / Захиалга / Сэтгэгдэл",
    "Хаяг, утас, нээх цаг харагдана",
    "QR код шаардлагагүй — хэн ч хандах",
    "Wouter routing, PublicLayout component",
    "Mobile-first responsive (320px–1920px)",
]
add_bullets(s, dx, dy + Inches(0.6), dw, Inches(4), features,
            size=12, bullet_color=PRIMARY)

footer(s, 8)


# ============================================================
# SLIDE 9 — Public Browse + Reservations
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "8. Public хуудас — Цэс ба ширээний урьдчилан захиалга")


def mock_browser(slide, x, y, w, h, title):
    add_round(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, w, Inches(0.35), RGBColor(0xE0, 0xE0, 0xE0))
    for i, c in enumerate([RGBColor(0xFF, 0x5F, 0x57),
                           RGBColor(0xFE, 0xBC, 0x2E),
                           RGBColor(0x28, 0xC8, 0x40)]):
        add_oval(slide, x + Inches(0.08 + i * 0.22),
                 y + Inches(0.08), Inches(0.15), Inches(0.15), c)
    add_text(slide, x + Inches(0.5), y + Inches(0.05),
             w - Inches(0.6), Inches(0.3),
             title, size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)


# LEFT: /browse
lx, ly, lw, lh = Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.7)
mock_browser(s, lx, ly, lw, lh, "/browse — Цэс үзэх")
add_round(s, lx + Inches(0.3), ly + Inches(0.6),
          lw - Inches(0.6), Inches(0.45), WHITE)
add_text(s, lx + Inches(0.5), ly + Inches(0.6),
         lw - Inches(1), Inches(0.45),
         "🔍 Хоол хайх...", size=11, color=MUTED,
         anchor=MSO_ANCHOR.MIDDLE)
chips = ["Бүгд", "Гол хоол", "Шөл", "Салат", "Уух зүйл"]
cx_curr = lx + Inches(0.3)
for i, ch in enumerate(chips):
    cw = Inches(0.95 if len(ch) > 4 else 0.7)
    fill = PRIMARY if i == 0 else WHITE
    color = WHITE if i == 0 else DARK
    add_round(s, cx_curr, ly + Inches(1.15), cw, Inches(0.3), fill,
              None if i == 0 else RGBColor(0xCC, 0xCC, 0xCC))
    add_text(s, cx_curr, ly + Inches(1.15), cw, Inches(0.3),
             ch, size=9, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx_curr += cw + Inches(0.08)

foods = [("Хуушуур", "5,500₮"), ("Бууз", "8,000₮"),
         ("Цуйван", "12,000₮"), ("Бөхлөг", "15,000₮")]
gx = lx + Inches(0.3)
gy = ly + Inches(1.7)
gw = (lw - Inches(0.6) - Inches(0.2)) / 2
gh = Inches(1.7)
for i, (name, price) in enumerate(foods):
    col = i % 2
    row = i // 2
    fx = gx + col * (gw + Inches(0.2))
    fy = gy + row * (gh + Inches(0.2))
    add_round(s, fx, fy, gw, gh, WHITE, RGBColor(0xCC, 0xCC, 0xCC))
    add_rect(s, fx, fy, gw, Inches(1.0), RGBColor(0xE5, 0xE7, 0xEB))
    add_text(s, fx, fy + Inches(0.3), gw, Inches(0.4),
             "🍴", size=20, align=PP_ALIGN.CENTER)
    add_text(s, fx + Inches(0.15), fy + Inches(1.05),
             gw - Inches(0.3), Inches(0.3),
             name, size=11, bold=True)
    add_text(s, fx + Inches(0.15), fy + Inches(1.35),
             gw - Inches(0.3), Inches(0.3),
             price, size=11, bold=True, color=PRIMARY)

# RIGHT: /reservations
rx, ry, rw, rh = Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.7)
mock_browser(s, rx, ry, rw, rh, "/reservations — Ширээ урьдчилан захиалах")
add_round(s, rx + Inches(0.3), ry + Inches(0.6),
          Inches(2.6), Inches(0.3), RGBColor(0xFF, 0xF4, 0xE0))
add_text(s, rx + Inches(0.3), ry + Inches(0.6),
         Inches(2.6), Inches(0.3),
         "📅 Урьдчилан захиалга", size=10, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, rx + Inches(0.3), ry + Inches(1.0),
         rw - Inches(0.6), Inches(0.5),
         "Захиалгаа өгөх", size=18, bold=True)

fields = [("Нэр", "Бат"), ("Утас", "9999-1234"),
          ("Хүний тоо", "4 хүн"),
          ("Огноо ба цаг", "2026-05-10  19:00"),
          ("Тэмдэглэл", "Цонхны дэргэд")]
fy_c = ry + Inches(1.6)
for label, val in fields:
    add_text(s, rx + Inches(0.3), fy_c,
             rw - Inches(0.6), Inches(0.2),
             label, size=9, color=MUTED)
    add_round(s, rx + Inches(0.3), fy_c + Inches(0.22),
              rw - Inches(0.6), Inches(0.32), WHITE,
              RGBColor(0xCC, 0xCC, 0xCC))
    add_text(s, rx + Inches(0.4), fy_c + Inches(0.22),
             rw - Inches(0.7), Inches(0.32),
             val, size=10, anchor=MSO_ANCHOR.MIDDLE)
    fy_c += Inches(0.6)
add_round(s, rx + Inches(0.3), fy_c,
          rw - Inches(0.6), Inches(0.4), PRIMARY)
add_text(s, rx + Inches(0.3), fy_c,
         rw - Inches(0.6), Inches(0.4),
         "Захиалга илгээх", size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 9)


# ============================================================
# SLIDE 10 — Public Reviews
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "9. Public хуудас — Сэтгэгдэл (Reviews)",
       "Дундаж үнэлгээ, түгээлт, шинэ сэтгэгдэл бичих")

sx, sy, sw_c, sh_c = Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.7)
add_round(s, sx, sy, sw_c, sh_c, LIGHT)
add_text(s, sx, sy + Inches(0.4), sw_c, Inches(0.6),
         "Зочдын үнэлгээ", size=22, bold=True, align=PP_ALIGN.CENTER)
add_text(s, sx, sy + Inches(1.1), sw_c, Inches(1.4),
         "⭐⭐⭐⭐⭐", size=72, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, sx, sy + Inches(2.7), sw_c, Inches(0.7),
         "4.8 / 5.0", size=42, bold=True, color=PRIMARY,
         align=PP_ALIGN.CENTER)
add_text(s, sx, sy + Inches(3.5), sw_c, Inches(0.5),
         "84 сэтгэгдэл", size=14, color=MUTED, align=PP_ALIGN.CENTER)

dist = [(5, 0.72), (4, 0.21), (3, 0.05), (2, 0.01), (1, 0.01)]
by = sy + Inches(4.1)
for star, frac in dist:
    add_text(s, sx + Inches(0.5), by, Inches(0.4), Inches(0.3),
             f"{star}★", size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, sx + Inches(1.0), by + Inches(0.07),
              sw_c - Inches(2.2), Inches(0.18), WHITE)
    add_round(s, sx + Inches(1.0), by + Inches(0.07),
              (sw_c - Inches(2.2)) * frac, Inches(0.18), ACCENT)
    add_text(s, sx + sw_c - Inches(0.9), by, Inches(0.6), Inches(0.3),
             f"{int(frac*84)}", size=11, color=MUTED,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    by += Inches(0.28)

rx, ry, rw, rh = Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.7)
add_round(s, rx, ry, rw, rh, LIGHT)
add_text(s, rx + Inches(0.3), ry + Inches(0.2),
         rw - Inches(0.6), Inches(0.4),
         "💬 Сүүлийн сэтгэгдэл", size=14, bold=True, color=PRIMARY)

reviews = [
    ("Бат-Эрдэнэ", 5, "Гайхалтай үйлчилгээ! Хоол маш амттай."),
    ("Сараа", 5, "QR кодоор захиалга хийсэн нь маш хурдан."),
    ("Мөнхбат", 4, "Ширээ урьдчилан захиалсан нь тохиромжтой."),
]
ry_c = ry + Inches(0.7)
for name, rating, comment in reviews:
    add_round(s, rx + Inches(0.3), ry_c,
              rw - Inches(0.6), Inches(0.95), WHITE)
    add_text(s, rx + Inches(0.4), ry_c + Inches(0.08),
             rw - Inches(0.8), Inches(0.3),
             "⭐" * rating, size=10, color=ACCENT)
    add_text(s, rx + Inches(0.4), ry_c + Inches(0.32),
             rw - Inches(0.8), Inches(0.25),
             name, size=10, bold=True)
    add_text(s, rx + Inches(0.4), ry_c + Inches(0.55),
             rw - Inches(0.8), Inches(0.4),
             f'"{comment}"', size=9, color=MUTED)
    ry_c += Inches(1.05)

add_round(s, rx + Inches(0.3), ry_c, rw - Inches(0.6), Inches(1.2), PRIMARY)
add_text(s, rx + Inches(0.5), ry_c + Inches(0.2),
         rw - Inches(1), Inches(0.3),
         "✏  Сэтгэгдэл үлдээх", size=12, bold=True, color=WHITE)
add_text(s, rx + Inches(0.5), ry_c + Inches(0.5),
         rw - Inches(1), Inches(0.3),
         "Үнэлгээ:  ⭐⭐⭐⭐⭐", size=11, color=WHITE)
add_round(s, rx + Inches(0.5), ry_c + Inches(0.85),
          Inches(1.5), Inches(0.3), WHITE)
add_text(s, rx + Inches(0.5), ry_c + Inches(0.85),
         Inches(1.5), Inches(0.3),
         "Илгээх", size=10, bold=True, color=PRIMARY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 10)


# ============================================================
# SLIDE 11 — Зочны мобайл цэс (with phone frame)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "10. Зочны мобайл интерфейс — Цэсийн хуудас",
       "QR кодоор нэвтэрсэн зочны утсан дээрх дэлгэц")

# Background gradient circle decoration
add_oval(s, Inches(-1), Inches(5), Inches(4), Inches(4),
         RGBColor(0xE8, 0xEF, 0xFA))

# Phone frame
phone_h = Inches(5.6)
phone_x = Inches(1.5)
phone_y = Inches(1.3)
add_phone_frame(s, FIG / "guest-menu.png", phone_x, phone_y, phone_h)

# Description right
dx = Inches(5.5)
dy = Inches(1.3)
add_text(s, dx, dy, Inches(7.5), Inches(0.5),
         "Гол элементүүд", size=20, bold=True, color=PRIMARY)
items = [
    "Толгой: рестораны нэр + ширээний дугаар",
    "Ангиллын таб (горизонтал scroll)",
    "Хоолны карт: зураг, нэр, тайлбар, үнэ",
    "Сагсанд нэмэх товч (бодит цагт сагс шинэчлэгдэнэ)",
    "Сагсны badge (онцлох тоо)",
    "Олон төхөөрөмжид нэг сагс синхрончлогдоно",
    "Mobile-first responsive дизайн",
]
add_bullets(s, dx, dy + Inches(0.7), Inches(7.5), Inches(4), items,
            size=13, bullet_color=PRIMARY)

# Tech details
add_round(s, dx, dy + Inches(4.6), Inches(7.5), Inches(1.3), LIGHT)
add_text(s, dx + Inches(0.3), dy + Inches(4.7),
         Inches(7), Inches(0.4),
         "Хэрэгжүүлэлт", size=13, bold=True, color=DARK)
add_text(s, dx + Inches(0.3), dy + Inches(5.05),
         Inches(7), Inches(0.8),
         "React 18 + TanStack Query + Tailwind CSS\n"
         "Socket.IO room: session_<qrToken>  ·  GET /api/menu/categories",
         size=11, color=MUTED)

footer(s, 11)


# ============================================================
# SLIDE 12 — Зочны сагс (with phone frame)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "11. Зочны мобайл интерфейс — Хамтын сагс",
       "Олон зочин нэг сагсыг хамтран бүрдүүлэх")

add_oval(s, Inches(-1), Inches(5), Inches(4), Inches(4),
         RGBColor(0xE8, 0xEF, 0xFA))

phone_h = Inches(5.6)
phone_x = Inches(1.5)
phone_y = Inches(1.3)
add_phone_frame(s, FIG / "guest-cart.png", phone_x, phone_y, phone_h)

dx = Inches(5.5)
dy = Inches(1.3)
add_text(s, dx, dy, Inches(7.5), Inches(0.5),
         "Хамтын сагс (Shared Cart)",
         size=20, bold=True, color=PRIMARY)
add_text(s, dx, dy + Inches(0.6), Inches(7.5), Inches(0.9),
         "Нэг ширээний хэд хэдэн зочин тус тусдаа утаснаасаа сагсыг хамтран "
         "бүрдүүлж чадна. Сервер тал shared_cart_items хүснэгтэд хадгалж, "
         "Socket.IO-аар real-time синхрончилдог.",
         size=12, color=DARK)

# Flow diagram
fy = dy + Inches(1.8)
flow = [
    ("📱 А зочин", "+ Хуушуур"),
    ("🌐 Server", "INSERT shared_cart_items"),
    ("📡 Socket.IO", "cart:updated → session room"),
    ("📱 Б, В, Г зочдын дэлгэц", "сагс шууд шинэчлэгдэнэ"),
]
for i, (a, b) in enumerate(flow):
    add_round(s, dx, fy + i * Inches(0.7),
              Inches(7.5), Inches(0.55), LIGHT)
    add_text(s, dx + Inches(0.2), fy + i * Inches(0.7),
             Inches(2.8), Inches(0.55),
             a, size=12, bold=True, color=PRIMARY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, dx + Inches(3.0), fy + i * Inches(0.7),
             Inches(4.3), Inches(0.55),
             b, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 12)


# ============================================================
# SLIDE 13 — Захиалгын статус (with phone frame)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "12. Зочны мобайл интерфейс — Захиалгын статус",
       "Бодит цагт шинэчлэгдэх timeline харуулалт")

add_oval(s, Inches(-1), Inches(5), Inches(4), Inches(4),
         RGBColor(0xE8, 0xEF, 0xFA))

phone_h = Inches(5.6)
phone_x = Inches(1.5)
phone_y = Inches(1.3)
add_phone_frame(s, FIG / "guest-order-status.png", phone_x, phone_y, phone_h)

dx = Inches(5.5)
dy = Inches(1.3)
add_text(s, dx, dy, Inches(7.5), Inches(0.5),
         "Захиалгын төлвийн шилжилт",
         size=20, bold=True, color=PRIMARY)

states = [
    ("pending", "Хүлээгдэж байна", PRIMARY, True),
    ("confirmed", "Баталгаажсан", PRIMARY, True),
    ("preparing", "Бэлтгэгдэж байна", ACCENT, True),
    ("ready", "Бэлэн", SUCCESS, False),
    ("served", "Хүргэгдсэн", MUTED, False),
    ("paid", "Төлөгдсөн", MUTED, False),
]
ty = dy + Inches(0.7)
for i, (key, label, color, active) in enumerate(states):
    cy = ty + i * Inches(0.65)
    add_oval(s, dx, cy + Inches(0.05),
             Inches(0.35), Inches(0.35),
             color if active else RGBColor(0xDD, 0xDD, 0xDD))
    if i < len(states) - 1:
        add_rect(s, dx + Inches(0.155), cy + Inches(0.4),
                 Inches(0.04), Inches(0.3),
                 color if active else RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s, dx + Inches(0.5), cy, Inches(2.5), Inches(0.4),
             key, size=11, bold=True,
             color=color if active else MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, dx + Inches(2.7), cy, Inches(4), Inches(0.4),
             label, size=11, color=DARK if active else MUTED,
             anchor=MSO_ANCHOR.MIDDLE)

footer(s, 13)


# ============================================================
# SLIDE 14 — QR архитектурын загвар
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "13. QR архитектурын загвар",
       "Stable identity + Ephemeral session JWT — хоёр давхаргат загвар")

# Two architectural pillars
left_x = Inches(0.6)
left_y = Inches(1.2)
left_w = Inches(6.0)
left_h = Inches(2.6)

add_round(s, left_x, left_y, left_w, left_h, RGBColor(0xE0, 0xE9, 0xF7))
add_rect(s, left_x, left_y, left_w, Inches(0.5), PRIMARY)
add_text(s, left_x, left_y, left_w, Inches(0.5),
         "Давхарга 1: Stable QR (тогтвортой ID)",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, left_x + Inches(0.3), left_y + Inches(0.7),
         left_w - Inches(0.6), Inches(0.4),
         "tables.qr_token  (UUID v4)",
         size=12, bold=True, color=DARK, font="Consolas")
add_bullets(s, left_x + Inches(0.4), left_y + Inches(1.15),
            left_w - Inches(0.8), Inches(2.0),
            [
                "Ширээ үүсэх үед нэг удаа автоматаар үүсгэгдэнэ",
                "Цаасан QR код-д хэвлэгдэж ширээнд наагдана",
                "Session дуусахад rotate ХИЙГДЭХГҮЙ — тогтвортой",
            ],
            size=11, bullet_color=PRIMARY)

right_x = Inches(6.8)
add_round(s, right_x, left_y, left_w, left_h, RGBColor(0xDC, 0xF2, 0xE5))
add_rect(s, right_x, left_y, left_w, Inches(0.5), SUCCESS)
add_text(s, right_x, left_y, left_w, Inches(0.5),
         "Давхарга 2: Ephemeral session JWT (хувь хүний security)",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, right_x + Inches(0.3), left_y + Inches(0.7),
         left_w - Inches(0.6), Inches(0.4),
         "JWT HS256  ·  payload: {sid, tid, cid}",
         size=12, bold=True, color=DARK, font="Consolas")
add_bullets(s, right_x + Inches(0.4), left_y + Inches(1.15),
            left_w - Inches(0.8), Inches(2.0),
            [
                "Зочин check-in хийхэд server тал гаргана (8 цаг expiry)",
                "localStorage-д хадгална, захиалгад X-Session-Token-аар явна",
                "Session хаагдвал JWT шалгалтад автомат хүчингүй болно",
            ],
            size=11, bullet_color=SUCCESS)

# Bottom: how they work together
by = Inches(4.05)
add_round(s, Inches(0.6), by, SW - Inches(1.2), Inches(2.9), LIGHT)
add_text(s, Inches(0.8), by + Inches(0.15),
         SW - Inches(1.6), Inches(0.5),
         "🛡  Аюулгүй байдлын онцлог: Photo-QR довтолгооноос хамгаалалт",
         size=14, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), by + Inches(0.65),
         SW - Inches(1.6), Inches(0.5),
         "QR код стабиль ч гэсэн зурагласан хэн ч захиалга хийж чадахгүй "
         "— дараах 3 давхарга хамгаалалт ажиллана:",
         size=11, color=DARK)

steps = [
    ("1", "QR scan үед",
     "Зөвхөн public цэс нээгдэнэ.\nЗахиалга хийх боломжгүй."),
    ("2", "Захиалга хийхэд",
     "Session JWT заавал шаардлагатай.\nStaff идэвхжүүлсний дараа л JWT олгогдоно."),
    ("3", "Session хаагдахад",
     "Хуучин JWT шалгалтад 401 invalid_session\nбуцаах тул дахин ашиглах боломжгүй."),
]
sx_step = Inches(0.8)
sy_step = by + Inches(1.3)
sw_step = (SW - Inches(2)) / 3
for i, (num, title, desc) in enumerate(steps):
    x = sx_step + i * (sw_step + Inches(0.1))
    add_round(s, x, sy_step, sw_step, Inches(1.45), WHITE,
              RGBColor(0xCC, 0xCC, 0xCC))
    add_oval(s, x + Inches(0.15), sy_step + Inches(0.15),
             Inches(0.4), Inches(0.4), PRIMARY)
    add_text(s, x + Inches(0.15), sy_step + Inches(0.15),
             Inches(0.4), Inches(0.4),
             num, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.7), sy_step + Inches(0.15),
             sw_step - Inches(0.85), Inches(0.4),
             title, size=12, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), sy_step + Inches(0.7),
             sw_step - Inches(0.4), Inches(0.7),
             desc, size=10, color=MUTED)

footer(s, 14)


# ============================================================
# SLIDE 15 — Multi-device + shared cart
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "14. Олон төхөөрөмжийн дэмжлэг",
       "session_participants junction + capacity-driven device limit")

# Center table
tx, ty = Inches(4.5), Inches(2.3)
tw, th = Inches(4.0), Inches(2.5)
add_round(s, tx, ty, tw, th, RGBColor(0x8B, 0x4F, 0x2A))
add_text(s, tx, ty + Inches(0.3), tw, Inches(0.6),
         "🍽 Ширээ #5", size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, tx, ty + Inches(0.95), tw, Inches(0.5),
         "capacity = 4", size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, tx, ty + Inches(1.45), tw, Inches(0.5),
         "qr_token (stable)", size=11,
         color=RGBColor(0xFF, 0xE0, 0xB0), align=PP_ALIGN.CENTER)

# 4 phone icons around
positions = [
    (Inches(2.0), Inches(1.5), "А зочин"),
    (Inches(9.5), Inches(1.5), "Б зочин"),
    (Inches(2.0), Inches(4.7), "В зочин"),
    (Inches(9.5), Inches(4.7), "Г зочин"),
]
for px, py, label in positions:
    add_round(s, px, py, Inches(1.7), Inches(2.0), LIGHT)
    add_round(s, px + Inches(0.3), py + Inches(0.2),
              Inches(1.1), Inches(1.4), DARK, adj=0.18)
    add_text(s, px + Inches(0.3), py + Inches(0.5),
             Inches(1.1), Inches(0.4),
             "📱", size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, px + Inches(0.3), py + Inches(0.95),
             Inches(1.1), Inches(0.5),
             "Сагс\n(хамтын)", size=8, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, px, py + Inches(1.6), Inches(1.7), Inches(0.4),
             label, size=11, bold=True, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom description
desc_y = Inches(5.3)
add_round(s, Inches(0.5), desc_y, Inches(12.3), Inches(1.6), LIGHT)
add_text(s, Inches(0.7), desc_y + Inches(0.1),
         Inches(12), Inches(0.4),
         "Ажиллах зарчим", size=14, bold=True, color=PRIMARY)
add_text(s, Inches(0.7), desc_y + Inches(0.5),
         Inches(12), Inches(0.4),
         "▸  Бүх 4 зочин QR-ыг скан хийж check-in → 4 өөр session JWT олгогдоно",
         size=12, color=DARK)
add_text(s, Inches(0.7), desc_y + Inches(0.85),
         Inches(12), Inches(0.4),
         "▸  Capacity-driven limit: 4 хүний ширээнд хамгийн ихдээ 4 төхөөрөмж "
         "(5 дахь нь 429 session_full)",
         size=12, color=DARK)
add_text(s, Inches(0.7), desc_y + Inches(1.2),
         Inches(12), Inches(0.4),
         "▸  Сагсанд нэмэх → cart:updated event → бүх 4 утсанд real-time нэг "
         "сагсаар синхрончлогдоно",
         size=12, color=DARK)

footer(s, 15)


# ============================================================
# SLIDE 16 — Use case (8-step guest flow)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "15. Use case — Зочны бүрэн урсгал")

steps = [
    ("1", "QR скан", "Зочин ширээний QR-ыг\nскан хийж цэс нээгдэнэ", "📱"),
    ("2", "Check-in",
     "\"Ширээнд сууж байна\" товч →\nServer session JWT олгоно", "🪑"),
    ("3", "Цэс үзэх",
     "Ангилал шилжих, хайх,\nдэлгэрэнгүй зураг харах", "🍽"),
    ("4", "Сагсанд нэмэх",
     "Бодит цагт хамтын сагсанд\nнэмэгдэх (multi-device sync)", "🛒"),
    ("5", "Захиалга өгөх",
     "POST /api/orders\nКассирт notification очино", "✅"),
    ("6", "Бэлтгэж байна",
     "Кассир баталгаажуулах →\nГал тогоо preparing → ready", "👨‍🍳"),
    ("7", "Хүлээн авах",
     "Кассир served гэж тэмдэглэх,\nstatus timeline шинэчлэгдэнэ", "🍴"),
    ("8", "Тооцоо + Сэтгэгдэл",
     "paid → session хаагдах →\nЗочин сэтгэгдэл үлдээх", "⭐"),
]
cw = Inches(3.0)
ch = Inches(2.7)
gap_x = Inches(0.18)
gap_y = Inches(0.18)
start_x = (SW - (cw * 4 + gap_x * 3)) / 2
start_y = Inches(1.2)
for i, (num, title, desc, ico) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = start_x + col * (cw + gap_x)
    y = start_y + row * (ch + gap_y)
    color = PRIMARY if row == 0 else SUCCESS
    add_round(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.5), ch, color)
    add_text(s, x, y + Inches(0.3), Inches(0.5), Inches(0.5),
             num, size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.6), y + Inches(0.15),
             cw - Inches(0.7), Inches(0.5),
             ico, size=22)
    add_text(s, x + Inches(0.6), y + Inches(0.85),
             cw - Inches(0.7), Inches(0.4),
             title, size=14, bold=True, color=DARK)
    add_text(s, x + Inches(0.6), y + Inches(1.3),
             cw - Inches(0.7), ch - Inches(1.4),
             desc, size=10, color=MUTED)

footer(s, 16)


# ============================================================
# SLIDE 17 — Staff dashboard (orders)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "16. Ажилтны самбар — Захиалгын удирдлага (Cashier view)")

img_path = FIG / "staff-orders.png"
add_image(s, img_path, Inches(0.5), Inches(1.2), w=Inches(8.5))

dx = Inches(9.3)
dy = Inches(1.3)
add_text(s, dx, dy, Inches(3.7), Inches(0.5),
         "Гол функцууд", size=18, bold=True, color=PRIMARY)
features = [
    "Захиалгууд card-аар",
    "Status filter: pending /\npreparing / ready / served",
    "Шинэ захиалга real-time\nанхааруулга",
    "Status шилжүүлэх (1 click)",
    "Ширээний дугаар + цаг +\nхоолын жагсаалт",
    "Notification bell — шинэ\nevent count",
    "Ээлж тус бүрийн дугаар\n(өглөө/үд/орой)",
]
add_bullets(s, dx, dy + Inches(0.6), Inches(3.7), Inches(5),
            features, size=11, bullet_color=PRIMARY)

footer(s, 17)


# ============================================================
# SLIDE 18 — Manager dashboard (menu + tables)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "17. Менежерийн самбар — Цэс ба ширээний удирдлага")

img1 = FIG / "staff-menu-manage.png"
img2 = FIG / "staff-tables.png"
add_image(s, img1, Inches(0.4), Inches(1.2), w=Inches(6.3))
add_image(s, img2, Inches(6.9), Inches(1.2), w=Inches(6.0))

add_text(s, Inches(0.4), Inches(6.5), Inches(6.3), Inches(0.4),
         "Цэс удирдлага: ангилал, хоол, үнэ, зураг",
         size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.9), Inches(6.5), Inches(6.0), Inches(0.4),
         "Ширээ удирдлага + QR код хэвлэх",
         size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

footer(s, 18)


# ============================================================
# SLIDE 19 — Reports
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "18. Тайлан, статистик")

img1 = FIG / "staff-report-stats.png"
img2 = FIG / "staff-report-charts.png"
add_image(s, img1, Inches(0.4), Inches(1.2), w=Inches(6.3))
add_image(s, img2, Inches(6.9), Inches(1.2), w=Inches(6.0))

add_round(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.0), LIGHT)
add_text(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.4),
         "Боломжит тайлангууд", size=13, bold=True, color=PRIMARY)
add_text(s, Inches(0.6), Inches(6.45), Inches(12), Inches(0.5),
         "▸ Өдрийн орлого  ·  ▸ Эрэлттэй хоол ТОП-10  ·  ▸ Ширээний ашиглалт  ·  "
         "▸ Доод хэмжээнд хүрсэн нөөц  ·  ▸ Ээлж тус бүрийн борлуулалт",
         size=11, color=DARK)

footer(s, 19)


# ============================================================
# SLIDE 20 — Туршилтын үр дүн ба дүгнэлт
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "19. Туршилтын үр дүн ба дүгнэлт")

# Top: 3 metric cards
my = Inches(1.15)
mh = Inches(1.6)
mc_w = (SW - Inches(1.0)) / 3 - Inches(0.15)

metrics = [
    ("📊", "SUS үнэлгээ", "84.2 / 100",
     "Маш сайн зэрэглэл, Grade A\n28 хэрэглэгчээр", SUCCESS),
    ("⚡", "k6 ачааллын туршилт", "120 мс / 45 мс",
     "API дундаж · WebSocket дундаж\n200 VU, 5 минут", PRIMARY),
    ("📈", "Уламжлалт vs QR", "−62% / −94%",
     "Захиалгын хугацаа · Алдааны хувь\n1.6 мин · 0.5%", ACCENT),
]
mx = Inches(0.5)
for i, (icon, title, value, sub, color) in enumerate(metrics):
    x = mx + i * (mc_w + Inches(0.15))
    add_round(s, x, my, mc_w, mh, LIGHT)
    add_rect(s, x, my, mc_w, Inches(0.05), color)
    add_text(s, x + Inches(0.2), my + Inches(0.15),
             Inches(0.6), Inches(0.5),
             icon, size=22)
    add_text(s, x + Inches(0.9), my + Inches(0.15),
             mc_w - Inches(1.1), Inches(0.4),
             title, size=12, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), my + Inches(0.6),
             mc_w - Inches(0.4), Inches(0.5),
             value, size=22, bold=True, color=DARK)
    add_text(s, x + Inches(0.2), my + Inches(1.05),
             mc_w - Inches(0.4), Inches(0.5),
             sub, size=10, color=MUTED)

# Bottom: achievements + future
ay_ = Inches(2.95)
add_round(s, Inches(0.5), ay_, Inches(6.2), Inches(4.0), LIGHT)
add_text(s, Inches(0.7), ay_ + Inches(0.15),
         Inches(5.8), Inches(0.5),
         "✅ Хэрэгжүүлсэн ажил", size=15, bold=True, color=SUCCESS)
done = [
    "11 хүснэгттэй PostgreSQL schema",
    "43 REST API endpoint, 12 Socket.IO event",
    "4 хэрэглэгчийн интерфейс",
    "QR хоёр давхарга: stable + ephemeral JWT",
    "Multi-device + shared cart real-time",
    "QR-гүй public хуудаснууд (4)",
    "k6 туршилт + SUS үнэлгээ 84.2",
    "Cloudflare Tunnel-аар online demo",
]
add_bullets(s, Inches(0.8), ay_ + Inches(0.7),
            Inches(5.6), Inches(3.2),
            done, size=11, bullet_color=SUCCESS, space_after=4)

fx = Inches(7.0)
add_round(s, fx, ay_, Inches(5.8), Inches(4.0),
          RGBColor(0xFD, 0xEB, 0xC9))
add_text(s, fx + Inches(0.3), ay_ + Inches(0.15),
         Inches(5.4), Inches(0.5),
         "🚀 Цаашид хөгжүүлэх боломж", size=15, bold=True, color=ACCENT)
future = [
    "QPay / SocialPay / Monpay цахим төлбөр",
    "AI-ээс хоол санал болгох",
    "Олон салбар рестораны модуль",
    "AR (Augmented Reality) 3D хоол",
    "React Native native апп",
    "Capacity planning / load forecasting",
    "Loyalty point / membership program",
    "Multilingual (EN, JP, KR)",
]
add_bullets(s, fx + Inches(0.3), ay_ + Inches(0.7),
            Inches(5.2), Inches(3.2),
            future, size=11, bullet_color=ACCENT, space_after=4)

footer(s, 20)


# ============================================================
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Slides: {len(prs.slides)}")
